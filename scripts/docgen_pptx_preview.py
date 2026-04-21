"""Split a PPTX into single-slide decks, then render each via qlmanage.

macOS-only — uses ``qlmanage`` as a lightweight slide previewer (it
only renders slide 1 of any deck, so we split and feed it one-by-one).

Usage:
    python3 scripts/docgen_pptx_preview.py <pptx_path>
"""

from __future__ import annotations

import shutil
import subprocess
import sys
from pathlib import Path


def split_pptx(src: Path, out_dir: Path) -> list[Path]:
    """Emit one single-slide .pptx per slide of ``src``.

    Uses python-pptx to copy the deck, then removes all slides except
    the i-th. NB python-pptx's "remove slide" helper doesn't ship;
    we hack the XML.
    """
    from pptx import Presentation

    src_prs = Presentation(str(src))
    n = len(src_prs.slides)
    out_paths: list[Path] = []
    for i in range(n):
        dst = out_dir / f"{src.stem}_slide_{i+1:02d}.pptx"
        shutil.copy2(src, dst)
        prs = Presentation(str(dst))
        # Remove slides that aren't i
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        for j, slide in enumerate(slides):
            if j != i:
                xml_slides.remove(slide)
        prs.save(str(dst))
        out_paths.append(dst)
    return out_paths


def render_thumbs(pptx_paths: list[Path], out_dir: Path, size: int = 1920) -> list[Path]:
    if not shutil.which("qlmanage"):
        print("qlmanage not found — macOS only", file=sys.stderr)
        return []
    out_dir.mkdir(parents=True, exist_ok=True)
    results: list[Path] = []
    for p in pptx_paths:
        subprocess.run(
            ["qlmanage", "-t", "-s", str(size), "-o", str(out_dir), str(p)],
            capture_output=True,
        )
        png = out_dir / f"{p.name}.png"
        if png.exists():
            results.append(png)
    return results


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: docgen_pptx_preview.py <pptx>", file=sys.stderr)
        return 1
    src = Path(sys.argv[1]).resolve()
    if not src.is_file():
        print(f"missing file: {src}", file=sys.stderr)
        return 1
    tmp = src.parent / ".preview_splits"
    if tmp.exists():
        shutil.rmtree(tmp)
    tmp.mkdir(parents=True)

    splits = split_pptx(src, tmp)
    out_dir = src.parent / f"{src.stem}_preview"
    out_dir.mkdir(parents=True, exist_ok=True)
    thumbs = render_thumbs(splits, out_dir)
    print(f"Rendered {len(thumbs)} slides → {out_dir}")
    for t in thumbs:
        print(f"  {t}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
