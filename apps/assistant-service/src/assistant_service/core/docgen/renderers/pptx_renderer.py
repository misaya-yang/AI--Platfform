"""PptxRenderer — design-system-driven dispatcher.

The renderer used to be a 2.3k-line class holding primitives + 14 layout
methods + block renderers + helpers. It was all split into:

* :mod:`primitives` — shape / text / chrome builders
* :mod:`layouts` — one module per cluster of layouts, plus a
  ``LAYOUT_DISPATCH`` dict mapping IR layout strings to free functions.
* :mod:`layouts.blocks` — block-level renderers shared across layouts
* :mod:`layouts.helpers` — text helpers (anchor / split / truncate /
  eyebrow extraction / table synthesis).

What lives here is only the public API contract: :class:`PptxRenderer`
with async ``render`` / ``fix``, plus the :class:`_Ctx` dataclass that
bundles design system + per-render metadata and gets threaded through
every layout.

Design conventions the layouts assume (distilled from 2026 SOTA research):

* Token-based colours (surface / ink / accent), not ``palette[0..4]``.
* Typography scale 1.333 (eyebrow / display / h1 / h2 / lead / body / caption).
* **Eyebrow kicker** above every H1.
* **Section numeral** behind content.
* **No accent underline under the title** — biggest AI-tell.
* **Asymmetric splits** (38/62), not 50/50.
* **Soft drop shadow** on every elevated surface via ``a:effectLst`` XML.
* **Gradient surface tint** for hero and card backgrounds.
* **Background geometry** — rotated soft rectangle at 4-8% fill.
"""

from __future__ import annotations

import logging
import re
import time
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from ..design_system import (
    ColorTokens,
    DesignSystem,
    TypeScale,
    available_systems,
    design_system_from_palette,
    get_design_system,
)
from ..ir import PptxIR
from .base import BaseRenderer, RenderError, RenderResult
from .layouts import LAYOUT_DISPATCH
from .primitives import SLIDE_H, SLIDE_W, Primitives

logger = logging.getLogger(__name__)


@dataclass
class _Ctx:
    """Per-render context — design system + typography + index."""

    ds: DesignSystem
    total_slides: int
    eyebrow_default: str
    brand: str
    author: str

    @property
    def c(self) -> ColorTokens:
        return self.ds.colors

    @property
    def t(self) -> TypeScale:
        return self.ds.type_scale


class PptxRenderer(BaseRenderer):
    format = "pptx"

    # ------------------------------------------------------------------ setup

    def _ctx(self, ir: PptxIR) -> _Ctx:
        # Pick design system: explicit metadata.extra["design_system"]
        # overrides; otherwise derive from the 5-colour palette bridge.
        ds_name = None
        if ir.metadata.model_extra:
            ds_name = ir.metadata.model_extra.get("design_system")
        ds: DesignSystem | None = None
        if ds_name and ds_name in available_systems():
            ds = get_design_system(ds_name)
        if ds is None:
            palette_hex = [c.value for c in ir.theme.palette]
            ds = design_system_from_palette(palette_hex)

        subtitle = ir.metadata.subtitle or ""
        eyebrow = subtitle.upper() if subtitle else ""
        brand = ir.metadata.author or "PRESENTATION"
        return _Ctx(
            ds=ds,
            total_slides=len(ir.content.slides),
            eyebrow_default=eyebrow,
            brand=brand,
            author=ir.metadata.author or "",
        )

    def _new_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        return prs

    # ------------------------------------------------------------------ main

    async def render(self, ir: PptxIR, out_dir: Path) -> RenderResult:
        if not isinstance(ir, PptxIR):
            raise RenderError(f"PptxRenderer got wrong IR type: {type(ir).__name__}")
        started = time.perf_counter()
        prs = self._new_presentation()
        ctx = self._ctx(ir)
        prims = Primitives()

        for idx, s_ir in enumerate(ir.content.slides, start=1):
            blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
            slide = prs.slides.add_slide(blank)
            fn = LAYOUT_DISPATCH.get(s_ir.layout, LAYOUT_DISPATCH["title_content"])
            fn(prs, slide, ir, s_ir, ctx, idx, prims)
            if s_ir.notes:
                slide.notes_slide.notes_text_frame.text = s_ir.notes

        out_dir.mkdir(parents=True, exist_ok=True)
        safe_name = re.sub(r"[^A-Za-z0-9_\-\.]", "_", ir.metadata.title)[:80] or "presentation"
        path = out_dir / f"{safe_name}.pptx"
        prs.save(str(path))

        return RenderResult(
            path=path,
            doc_type="pptx",
            bytes_size=path.stat().st_size,
            duration_ms=int((time.perf_counter() - started) * 1000),
            extra={"slides": len(ir.content.slides), "design_system": ctx.ds.name},
        )

    async def fix(self, ir: PptxIR, critic_findings, out_dir: Path) -> RenderResult:
        del critic_findings
        return await self.render(ir, out_dir)
