"""Low-level shape & text primitives for the PPTX renderer.

Split out of ``pptx_renderer.py`` so the 14 layout modules can share one
coherent vocabulary (``rect / ellipse / text / eyebrow / footer``) without
every layout re-importing ``pptx`` internals.

A single :class:`Primitives` instance is created per render and passed as
the last positional arg to every layout free function. It holds no
per-slide state — only the slide geometry constants and a logger.
"""

from __future__ import annotations

import logging
from typing import Optional

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ..design_system import tint
from .pptx_effects import (
    add_outer_shadow,
    set_linear_gradient_fill,
    set_no_stroke,
    set_rotation,
    set_transparency,
)

logger = logging.getLogger(__name__)


SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.66          # outer safe area (≈ 48pt)
GRID_GAP = 0.33


def rgb(hex_value: str) -> RGBColor:
    return RGBColor(
        int(hex_value[0:2], 16),
        int(hex_value[2:4], 16),
        int(hex_value[4:6], 16),
    )


class Primitives:
    """Shared shape / text builders used by every layout module.

    Stateless — the same instance is reused for every slide.
    """

    # ----------------------------------------------------------- shapes

    def rect(self, slide, *, left, top, width, height, hex_fill, shadow=False, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(hex_fill)
        set_no_stroke(shp)
        if radius:
            try:
                shp.adjustments[0] = 0.05
            except (IndexError, AttributeError) as exc:
                logger.debug("rounded-rect adjustments unavailable: %s", exc)
        if shadow:
            add_outer_shadow(shp, blur_emu=40_000, dist_emu=23_000, alpha_per_mille=18_000)
        return shp

    def gradient_rect(self, slide, *, left, top, width, height, stops, angle=90.0, shadow=False):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        set_no_stroke(shp)
        set_linear_gradient_fill(shp, stops=stops, angle_deg=angle)
        if shadow:
            add_outer_shadow(shp, blur_emu=40_000, dist_emu=23_000, alpha_per_mille=18_000)
        return shp

    def ellipse(self, slide, *, left, top, width, height, hex_fill, shadow=False):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = rgb(hex_fill)
        set_no_stroke(shp)
        if shadow:
            add_outer_shadow(shp, blur_emu=30_000, dist_emu=15_000, alpha_per_mille=15_000)
        return shp

    # ----------------------------------------------------------- text

    def text(
        self,
        slide,
        text: str,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        size_pt: float,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        v_anchor: str = "top",
        rgb_hex: str = "0F172A",
        font: str = "Calibri",
        line_spacing: Optional[float] = None,
        tracking_pct: float = 0.0,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        # CRITICAL: disable auto-resize so text boxes never grow beyond
        # their declared rectangle. LibreOffice respects this; PowerPoint
        # respects it. Without it, long text silently overflows into
        # neighbouring cards / the footer.
        try:
            tf.auto_size = MSO_AUTO_SIZE.NONE
        except (AttributeError, ValueError) as exc:
            logger.debug("auto_size=NONE not available on this textframe: %s", exc)
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(v_anchor, MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(align, PP_ALIGN.LEFT)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = rgb(rgb_hex)
        # Letter tracking via XML (a:rPr spc attribute, 100ths of point)
        if tracking_pct != 0.0:
            try:
                _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                rpr = run._r.find(f"{{{_A}}}rPr")
                if rpr is not None:
                    spc_units = int(tracking_pct * 1000)  # ~em thousandths
                    rpr.set("spc", str(spc_units))
            except (AttributeError, ValueError) as exc:
                logger.debug("letter-tracking XML injection failed: %s", exc)
        return box

    # ----------------------------------------------------------- chrome

    def background(self, slide, hex_fill: str) -> None:
        bg = self.rect(slide, left=0, top=0, width=SLIDE_W, height=SLIDE_H, hex_fill=hex_fill)
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def background_gradient(self, slide, from_hex: str, to_hex: str, angle: float = 135.0) -> None:
        bg = self.gradient_rect(
            slide, left=0, top=0, width=SLIDE_W, height=SLIDE_H,
            stops=[(0.0, from_hex), (1.0, to_hex)], angle=angle,
        )
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def eyebrow(self, slide, text_value: str, *, left: float, top: float, ctx, width: float = 8.0):
        """Small uppercase, tracked text above an H1. The single biggest
        "looks designed" lever."""
        text_value = text_value.strip()
        if not text_value:
            return None
        # Small coloured square + label
        dot_size = 0.14
        self.rect(
            slide,
            left=left,
            top=top + 0.06,
            width=dot_size,
            height=dot_size,
            hex_fill=ctx.c.accent,
        )
        return self.text(
            slide,
            text_value.upper(),
            left=left + dot_size + 0.18,
            top=top,
            width=width,
            height=0.3,
            size_pt=ctx.t.eyebrow_pt,
            bold=True,
            rgb_hex=ctx.c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=ctx.t.eyebrow_tracking_pct,
        )

    def section_numeral(self, slide, index: int, *, ctx, x: float, y: float) -> None:
        """Giant muted numeral in the corner — magazine-style decoration."""
        label = f"{index:02d} / {ctx.total_slides:02d}"
        self.text(
            slide,
            label,
            left=x,
            top=y,
            width=4.5,
            height=2.5,
            size_pt=ctx.t.section_numeral_pt * 0.6,
            bold=True,
            rgb_hex=tint(ctx.c.ink_muted, 0.6),
            font=ctx.ds.font_display,
            line_spacing=1.0,
            v_anchor="top",
        )

    def footer(self, slide, *, ctx, index: int) -> None:
        """Small footer bar with brand + page number."""
        self.text(
            slide,
            (ctx.brand or "").upper(),
            left=MARGIN,
            top=SLIDE_H - 0.42,
            width=SLIDE_W / 2,
            height=0.28,
            size_pt=9,
            bold=True,
            rgb_hex=ctx.c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=0.08,
        )
        self.text(
            slide,
            f"{index:02d} / {ctx.total_slides:02d}",
            left=SLIDE_W - MARGIN - 2.0,
            top=SLIDE_H - 0.42,
            width=2.0,
            height=0.28,
            size_pt=9,
            bold=True,
            align="right",
            rgb_hex=ctx.c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=0.08,
        )

    def background_geometry(self, slide, *, ctx) -> None:
        """Subtle decoration — rotated faint square + corner quarter-circle."""
        sq = self.rect(
            slide,
            left=-1.5,
            top=SLIDE_H - 2.2,
            width=3.5,
            height=3.5,
            hex_fill=ctx.c.accent,
        )
        set_transparency(sq, alpha_per_mille=6_000)
        set_rotation(sq, degrees=15.0)

        circ = self.ellipse(
            slide,
            left=SLIDE_W - 2.0,
            top=-2.0,
            width=4.0,
            height=4.0,
            hex_fill=ctx.c.accent_secondary,
        )
        set_transparency(circ, alpha_per_mille=5_000)

    def dot_grid(
        self,
        slide,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        ctx,
        dot_color: str,
        spacing: float = 0.28,
        dot_size: float = 0.04,
    ) -> None:
        """Low-contrast dot pattern for background texture."""
        rows = int(height / spacing)
        cols = int(width / spacing)
        for r in range(rows):
            for cx in range(cols):
                self.ellipse(
                    slide,
                    left=left + cx * spacing,
                    top=top + r * spacing,
                    width=dot_size,
                    height=dot_size,
                    hex_fill=dot_color,
                )
