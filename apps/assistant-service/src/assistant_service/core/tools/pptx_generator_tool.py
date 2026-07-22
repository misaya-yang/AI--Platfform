"""
PowerPoint Generation Tool for Assistant Service

Provides PowerPoint (.pptx) generation capabilities using python-pptx.
Creates professional presentations from structured JSON input.

Features:
- Title slides with subtitle
- Content slides with bullet points
- Two-column layout slides
- Image slides (with placeholder support)
- Professional templates with consistent styling
"""

from __future__ import annotations

import base64
import io
import re
import time
from dataclasses import dataclass
from typing import Any

from ai_gateway_core.logging import get_logger

from .tool_registry import (
    ToolCallRequest,
    ToolCallResult,
    ToolCategory,
    ToolDefinition,
    ToolExecutor,
    ToolParameter,
    ToolRiskLevel,
    register_tool,
)

logger = get_logger(__name__)


# =============================================================================
# PPT Generation Tool Definition
# =============================================================================

PPTX_GENERATION_DEFINITION = ToolDefinition(
    name="generate_pptx",
    description="""Generate a PowerPoint presentation (.pptx) file.

**WORKFLOW:**
1. You have ALREADY planned the content in the chat.
2. Now, simply CALL this tool with that content as structured JSON.
3. Do not output any more text. Just call the tool.

**REQUIRED PARAMETERS:**
- title: string
- slides: array of objects (see schema)
""",
    parameters=[
        ToolParameter(
            name="title",
            type="string",
            description="Presentation title (becomes filename and title slide).",
            required=True,
        ),
        ToolParameter(
            name="subtitle",
            type="string",
            description="Subtitle for title slide (author, date, organization).",
            required=False,
            default="",
        ),
        ToolParameter(
            name="slides",
            type="array",
            description="Array of slide objects. DO NOT output this in chat - it's internal only.",
            required=True,
            items={
                "type": "object",
                "properties": {
                    "type": {
                        "type": "string",
                        "description": "Slide type",
                        "enum": ["title", "content", "two_column", "section", "blank"],
                    },
                    "title": {"type": "string", "description": "Slide title"},
                    "subtitle": {"type": "string", "description": "Subtitle for section slides"},
                    "bullets": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Bullet points",
                    },
                    "left_content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Left column content",
                    },
                    "right_content": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Right column content",
                    },
                },
                "required": ["type"],
            },
        ),
        ToolParameter(
            name="theme",
            type="string",
            description="Color theme for the presentation.",
            required=False,
            default="professional",
            enum=["professional", "modern", "elegant", "vibrant"],
        ),
    ],
    category=ToolCategory.GENERATION,
    risk_level=ToolRiskLevel.LOW,
    capability_metadata={"operation_kind": "write"},
    when_to_use="""Use when user asks for PowerPoint/PPT/slides/演示文稿.
You MUST call this tool in the same response - do not just describe what you will do.""",
    when_not_to_use="Don't use for simple text. Don't use for documents (use Word instead).",
    examples=[],
    timeout_seconds=60,
)


# =============================================================================
# Theme Configurations
# =============================================================================


@dataclass
class ThemeConfig:
    """Color theme configuration for presentations."""

    primary_color: str  # RGB hex for primary elements
    secondary_color: str  # RGB hex for secondary elements
    accent_color: str  # RGB hex for accents
    text_color: str  # RGB hex for main text
    background_color: str  # RGB hex for background
    title_font: str  # Font name for titles
    body_font: str  # Font name for body text


THEMES = {
    "professional": ThemeConfig(
        primary_color="1F4E79",  # Dark blue
        secondary_color="2E75B6",  # Medium blue
        accent_color="5B9BD5",  # Light blue
        text_color="2F2F2F",  # Dark gray
        background_color="FFFFFF",  # White
        title_font="Calibri Light",
        body_font="Calibri",
    ),
    "modern": ThemeConfig(
        primary_color="00796B",  # Teal
        secondary_color="26A69A",  # Light teal
        accent_color="80CBC4",  # Lighter teal
        text_color="263238",  # Blue-gray
        background_color="FFFFFF",  # White
        title_font="Segoe UI Light",
        body_font="Segoe UI",
    ),
    "elegant": ThemeConfig(
        primary_color="37474F",  # Dark blue-gray
        secondary_color="546E7A",  # Medium blue-gray
        accent_color="78909C",  # Light blue-gray
        text_color="ECEFF1",  # Light gray (for dark bg)
        background_color="263238",  # Dark background
        title_font="Georgia",
        body_font="Georgia",
    ),
    "vibrant": ThemeConfig(
        primary_color="E65100",  # Deep orange
        secondary_color="FF6D00",  # Orange
        accent_color="FFAB40",  # Light orange
        text_color="212121",  # Almost black
        background_color="FFFFFF",  # White
        title_font="Arial Black",
        body_font="Arial",
    ),
}


# =============================================================================
# PowerPoint Generator
# =============================================================================


class PPTXGenerator:
    """Generate PowerPoint presentations using python-pptx."""

    def __init__(self):
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
            from pptx.util import Inches, Pt

            _ = Presentation, RGBColor, MSO_ANCHOR, PP_ALIGN, Inches, Pt
            self._pptx_available = True
        except ImportError:
            self._pptx_available = False
            logger.warning("python-pptx not installed, PPT generation disabled")

    @property
    def is_available(self) -> bool:
        return self._pptx_available

    def generate(
        self,
        title: str,
        subtitle: str,
        slides: list[dict[str, Any]],
        theme: str = "professional",
    ) -> bytes:
        """Generate a PowerPoint presentation.

        Args:
            title: Presentation title
            subtitle: Subtitle for title slide
            slides: List of slide definitions
            theme: Color theme name

        Returns:
            bytes: PowerPoint file content
        """
        if not self._pptx_available:
            raise RuntimeError("python-pptx is not installed")

        from pptx import Presentation
        from pptx.util import Inches

        # Get theme configuration
        theme_config = THEMES.get(theme, THEMES["professional"])

        # Create presentation with widescreen aspect ratio
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Add title slide
        self._add_title_slide(prs, title, subtitle, theme_config)

        # Add content slides
        for slide_def in slides:
            slide_type = slide_def.get("type", "content")

            if slide_type == "title":
                # Additional title slide
                self._add_title_slide(
                    prs, slide_def.get("title", ""), slide_def.get("subtitle", ""), theme_config
                )
            elif slide_type == "section":
                self._add_section_slide(
                    prs, slide_def.get("title", ""), slide_def.get("subtitle", ""), theme_config
                )
            elif slide_type == "content":
                self._add_content_slide(
                    prs, slide_def.get("title", ""), slide_def.get("bullets", []), theme_config
                )
            elif slide_type == "two_column":
                self._add_two_column_slide(
                    prs,
                    slide_def.get("title", ""),
                    slide_def.get("left_content", []),
                    slide_def.get("right_content", []),
                    theme_config,
                )
            elif slide_type == "blank":
                self._add_blank_slide(prs, theme_config)

        # Save to bytes
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        return buffer.read()

    def _hex_to_rgb(self, hex_color: str):
        """Convert hex color to RGBColor."""
        from pptx.dml.color import RGBColor

        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)
        return RGBColor(r, g, b)

    def _add_title_slide(self, prs, title: str, subtitle: str, theme: ThemeConfig):
        """Add a title slide."""
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        # Use blank layout and create custom design
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)

        # Add colored header band
        shapes = slide.shapes
        header = shapes.add_shape(
            1,  # Rectangle
            Inches(0),
            Inches(0),
            prs.slide_width,
            Inches(2.5),
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self._hex_to_rgb(theme.primary_color)
        header.line.fill.background()

        # Add title text
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.333), Inches(1.2))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb("FFFFFF")
        title_para.alignment = PP_ALIGN.CENTER

        # Add subtitle
        if subtitle:
            subtitle_box = shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(12.333), Inches(0.8))
            sub_frame = subtitle_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle
            sub_para.font.size = Pt(24)
            sub_para.font.color.rgb = self._hex_to_rgb(theme.secondary_color)
            sub_para.alignment = PP_ALIGN.CENTER

    def _add_section_slide(self, prs, title: str, subtitle: str, theme: ThemeConfig):
        """Add a section header slide."""
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes

        # Full-width colored band in the middle
        band = shapes.add_shape(1, Inches(0), Inches(2.5), prs.slide_width, Inches(2.5))
        band.fill.solid()
        band.fill.fore_color.rgb = self._hex_to_rgb(theme.secondary_color)
        band.line.fill.background()

        # Section title
        title_box = shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1.2))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(40)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb("FFFFFF")
        title_para.alignment = PP_ALIGN.CENTER

        # Section subtitle
        if subtitle:
            sub_box = shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.6))
            sub_frame = sub_box.text_frame
            sub_para = sub_frame.paragraphs[0]
            sub_para.text = subtitle
            sub_para.font.size = Pt(20)
            sub_para.font.italic = True
            sub_para.font.color.rgb = self._hex_to_rgb("EEEEEE")
            sub_para.alignment = PP_ALIGN.CENTER

    def _add_content_slide(self, prs, title: str, bullets: list[str], theme: ThemeConfig):
        """Add a content slide with title and bullet points."""
        from pptx.util import Inches, Pt

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes

        # Title bar
        title_bar = shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(theme.primary_color)
        title_bar.line.fill.background()

        # Slide title
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb("FFFFFF")

        # Content area with bullets
        if bullets:
            content_box = shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.733), Inches(5.4))
            tf = content_box.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

                p.text = bullet
                p.font.size = Pt(22)
                p.font.color.rgb = self._hex_to_rgb(theme.text_color)
                p.level = 0
                p.space_before = Pt(12)
                p.space_after = Pt(6)

                # Add bullet character
                p.bullet = True

    def _add_two_column_slide(
        self, prs, title: str, left_content: list[str], right_content: list[str], theme: ThemeConfig
    ):
        """Add a two-column slide."""
        from pptx.util import Inches, Pt

        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        shapes = slide.shapes

        # Title bar
        title_bar = shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.2))
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self._hex_to_rgb(theme.primary_color)
        title_bar.line.fill.background()

        # Slide title
        title_box = shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.7))
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = self._hex_to_rgb("FFFFFF")

        # Left column
        left_box = shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.4))
        self._fill_column(left_box, left_content, theme)

        # Right column
        right_box = shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(5.4))
        self._fill_column(right_box, right_content, theme)

        # Vertical divider
        divider = shapes.add_shape(1, Inches(6.5), Inches(1.6), Inches(0.03), Inches(5.4))
        divider.fill.solid()
        divider.fill.fore_color.rgb = self._hex_to_rgb(theme.accent_color)
        divider.line.fill.background()

    def _fill_column(self, textbox, content: list[str], theme: ThemeConfig):
        """Fill a column textbox with bullet content."""
        from pptx.util import Pt

        tf = textbox.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = self._hex_to_rgb(theme.text_color)
            p.space_before = Pt(8)
            p.space_after = Pt(4)
            p.bullet = True

    def _add_blank_slide(self, prs, theme: ThemeConfig):
        """Add a blank slide."""
        del theme
        blank_layout = prs.slide_layouts[6]
        prs.slides.add_slide(blank_layout)


# =============================================================================
# PPTX Generator Tool Executor
# =============================================================================


class PPTXGeneratorExecutor(ToolExecutor):
    """Executor for PowerPoint generation tool."""

    def __init__(self):
        self.generator = PPTXGenerator()

    async def execute(self, request: ToolCallRequest) -> ToolCallResult:
        """Execute PowerPoint generation."""
        start_time = time.time()

        title = request.arguments.get("title", "Presentation")
        subtitle = request.arguments.get("subtitle", "")
        slides = request.arguments.get("slides", [])
        theme = request.arguments.get("theme", "professional")

        if not slides:
            return ToolCallResult(
                call_id=request.call_id,
                tool_name=request.tool_name,
                success=False,
                error="Slides array is required and cannot be empty",
            )

        # Validate slides structure
        if not isinstance(slides, list):
            return ToolCallResult(
                call_id=request.call_id,
                tool_name=request.tool_name,
                success=False,
                error="Slides must be an array of slide objects",
            )

        try:
            pptx_bytes = self.generator.generate(
                title=title,
                subtitle=subtitle,
                slides=slides,
                theme=theme,
            )

            filename = f"{self._sanitize_filename(title)}.pptx"
            mime_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            duration_ms = (time.time() - start_time) * 1000

            return ToolCallResult(
                call_id=request.call_id,
                tool_name=request.tool_name,
                success=True,
                result=f"PowerPoint presentation '{filename}' generated successfully with {len(slides) + 1} slides (including title slide). The user can download it directly from the chat interface.",
                output_files=[
                    {
                        "filename": filename,
                        "content_base64": base64.b64encode(pptx_bytes).decode("utf-8"),
                        "mime_type": mime_type,
                        "size_bytes": len(pptx_bytes),
                    }
                ],
                metadata={
                    "format": "pptx",
                    "title": title,
                    "filename": filename,
                    "slide_count": len(slides) + 1,
                    "theme": theme,
                    "size_bytes": len(pptx_bytes),
                    "duration_ms": duration_ms,
                },
                duration_ms=duration_ms,
            )

        except Exception as e:
            logger.error(f"PPTX generation failed: {e}")
            return ToolCallResult(
                call_id=request.call_id,
                tool_name=request.tool_name,
                success=False,
                error=str(e),
                duration_ms=(time.time() - start_time) * 1000,
            )

    def _sanitize_filename(self, title: str) -> str:
        """Sanitize title for use as filename."""
        safe = re.sub(r'[<>:"/\\|?*]', "_", title)
        return safe[:100]


# =============================================================================
# Registration Helper
# =============================================================================


def register_pptx_generation_tool() -> bool:
    """Register PowerPoint generation tool with the global registry."""
    executor = PPTXGeneratorExecutor()

    if not executor.generator.is_available:
        logger.warning("python-pptx not installed, PPTX generation tool not registered")
        return False

    register_tool(PPTX_GENERATION_DEFINITION, executor)
    logger.info("Registered PowerPoint generation tool")
    return True
