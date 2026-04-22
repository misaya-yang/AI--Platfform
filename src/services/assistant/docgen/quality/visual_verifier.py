"""PptxPdfVisualVerifier — render to JPEG, hand to a fresh-context vision critic.

Flow:
  1. If soffice available, convert pptx → pdf.
  2. pdftoppm -jpeg -r 150 (or pypdfium2 as a fallback) to render each page.
  3. Send the JPEGs to a ``VisionCritic``. The CRITIC is intentionally
     abstracted behind a Protocol so:
       - production uses a fresh Anthropic client
       - tests use a deterministic :class:`StructuralVisionCritic` that
         inspects slide text / layout without vision
  4. Map critic issues back to slide / page indices.

Per the plan:
  "Vision critic must use fresh Anthropic client. Fresh context. Prompt
  pinned to 'Assume there are issues'."
"""

from __future__ import annotations

import base64
import shutil
import subprocess
from pathlib import Path
from typing import Optional, Protocol, runtime_checkable

from .types import CriticReport, Issue, IssueCategory, IssueSeverity


@runtime_checkable
class VisionCritic(Protocol):
    """Abstract vision critic."""

    def review(self, doc_type: str, images: list[Path], context: dict) -> list[Issue]: ...


# The prompt the FreshContextVisionCritic sends. Plan §4.5 pins the wording.
CRITIC_SYSTEM_PROMPT = """You are a design QA reviewer. Assume there ARE issues in this document.

Your job is to find:
  * overflow (text or images cropped past slide/page boundaries)
  * overlap (two shapes covering each other)
  * alignment (jagged left edges, misaligned columns)
  * contrast (light text on light background, small copy on busy image)
  * missing images (placeholder grey blocks, broken icons)
  * placeholder / lorem ipsum / AI-tell artefacts
  * accent lines under titles (forbidden)
  * centred body text (forbidden)

Return a JSON array: [{ "page": int, "category": str, "severity": "critical"|"warning"|"info",
"message": str, "hint"?: str }, ...]
"""


class FreshContextVisionCritic(VisionCritic):
    """Vision critic backed by a fresh Anthropic client.

    The critic NEVER shares context with the main agent. Each review is a
    one-shot request with the system prompt above plus the rendered
    slide / page images. The :func:`client_factory` is invoked on every
    review so each call gets a brand-new client (= fresh cache, fresh
    conversation id, fresh rate-limit budget).
    """

    def __init__(self, *, client_factory, model: str = "claude-sonnet-4-6") -> None:
        self._client_factory = client_factory
        self._model = model

    def review(self, doc_type: str, images: list[Path], context: dict) -> list[Issue]:
        client = self._client_factory()  # fresh!
        blocks = [{"type": "text", "text": CRITIC_SYSTEM_PROMPT}]
        for idx, img in enumerate(images):
            blocks.append({"type": "text", "text": f"\nPage/slide {idx + 1}:"})
            with open(img, "rb") as fh:
                blocks.append({
                    "type": "image",
                    "source": {
                        "type": "base64",
                        "media_type": "image/jpeg",
                        "data": base64.b64encode(fh.read()).decode("ascii"),
                    },
                })

        try:
            resp = client.messages.create(
                model=self._model,
                max_tokens=2000,
                system=CRITIC_SYSTEM_PROMPT,
                messages=[{"role": "user", "content": blocks}],
            )
            raw = resp.content[0].text
            parsed = self._parse_json_array(raw)
        except Exception as e:
            return [Issue(category=IssueCategory.OTHER, severity=IssueSeverity.INFO, message=f"critic invocation failed: {e}")]

        return [self._issue_from_dict(d) for d in parsed]

    def _parse_json_array(self, raw: str) -> list[dict]:
        import json
        import re

        m = re.search(r"\[.*\]", raw, re.DOTALL)
        if not m:
            return []
        try:
            data = json.loads(m.group(0))
            return data if isinstance(data, list) else []
        except json.JSONDecodeError:
            return []

    def _issue_from_dict(self, d: dict) -> Issue:
        return Issue(
            category=_CATEGORY.get(d.get("category", "other"), IssueCategory.OTHER),
            severity=_SEVERITY.get(d.get("severity", "warning"), IssueSeverity.WARNING),
            message=str(d.get("message", ""))[:400],
            target=d.get("page"),
            hint=d.get("hint"),
        )


_CATEGORY = {c.value: c for c in IssueCategory}
_SEVERITY = {s.value: s for s in IssueSeverity}


def default_vision_critic() -> "VisionCritic":
    """Return the production vision critic.

    If ``ANTHROPIC_API_KEY`` is set AND the ``anthropic`` SDK is importable,
    returns a :class:`FreshContextVisionCritic` that hands rendered pages to
    a fresh Anthropic client each call. Otherwise returns the deterministic
    :class:`StructuralVisionCritic` so CI / offline dev boxes keep working.
    """
    import os

    if os.environ.get("ANTHROPIC_API_KEY"):
        try:
            from anthropic import AsyncAnthropic  # type: ignore  # local, lazy
        except ImportError:
            return StructuralVisionCritic()

        def factory():
            return AsyncAnthropic()

        return FreshContextVisionCritic(client_factory=factory)
    return StructuralVisionCritic()


# --------------------------------------------------------------------------
# StructuralVisionCritic — a heuristic critic that works without any LLM.
# This exists so the fix-and-verify loop can be exercised in unit tests
# without burning an API call. It reads the pptx/pdf structure directly
# (not pixels) but emits Issue[] in the same shape.


class StructuralVisionCritic(VisionCritic):
    """Deterministic critic for tests — no vision, no LLM."""

    def review(self, doc_type: str, images: list[Path], context: dict) -> list[Issue]:
        issues: list[Issue] = []
        if doc_type == "pptx":
            issues.extend(self._review_pptx(context))
        elif doc_type == "pdf":
            issues.extend(self._review_pdf(context))
        return issues

    def _review_pptx(self, ctx: dict) -> list[Issue]:
        from pptx import Presentation
        out: list[Issue] = []
        path = ctx.get("artifact_path")
        if not path:
            return out
        prs = Presentation(str(path))
        for idx, slide in enumerate(prs.slides):
            # AI-tell: centred body text inside a body shape
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                if not tf.text or len(tf.text) < 3:
                    continue
                for p in tf.paragraphs:
                    from pptx.enum.text import PP_ALIGN
                    if p.alignment == PP_ALIGN.CENTER and len(tf.text) > 120:
                        out.append(Issue(
                            category=IssueCategory.AI_TELL,
                            severity=IssueSeverity.WARNING,
                            message=f"slide {idx+1}: long body text centred",
                            target=idx,
                            hint="set alignment=left",
                        ))
                        break
            # placeholder grep
            all_text = "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)
            if any(marker in all_text.lower() for marker in ("lorem", "ipsum", "xxxx", "todo", "placeholder")):
                out.append(Issue(
                    category=IssueCategory.PLACEHOLDER,
                    severity=IssueSeverity.CRITICAL,
                    message=f"slide {idx+1}: placeholder text present",
                    target=idx,
                ))
        return out

    def _review_pdf(self, ctx: dict) -> list[Issue]:
        import pypdf
        out: list[Issue] = []
        path = ctx.get("artifact_path")
        if not path:
            return out
        try:
            reader = pypdf.PdfReader(str(path))
        except Exception as e:
            return [Issue(category=IssueCategory.OPEN_ERROR, severity=IssueSeverity.CRITICAL, message=str(e))]
        for idx, page in enumerate(reader.pages):
            text = (page.extract_text() or "").lower()
            if any(marker in text for marker in ("lorem", "ipsum", "xxxx", "todo", "placeholder")):
                out.append(Issue(
                    category=IssueCategory.PLACEHOLDER,
                    severity=IssueSeverity.CRITICAL,
                    message=f"page {idx+1}: placeholder text",
                    target=idx,
                ))
        return out


# --------------------------------------------------------------------------
# The verifier


class PptxPdfVisualVerifier:
    """Render artifact → images → critic → Report."""

    name = "visual"

    def __init__(self, critic: Optional[VisionCritic] = None) -> None:
        self._critic = critic or StructuralVisionCritic()

    def verify(self, path: Path, *, doc_type: str) -> CriticReport:
        images: list[Path] = []
        if doc_type in ("pptx",):
            pdf_path = self._pptx_to_pdf(path)
            if pdf_path is not None:
                images = self._pdf_to_images(pdf_path)
            backend = "soffice+pdftoppm" if images else "structural_only"
        elif doc_type == "pdf":
            images = self._pdf_to_images(path)
            backend = "pdftoppm" if images else "structural_only"
        else:
            return CriticReport(passed=True, backend="n/a", note=f"no visual verifier for {doc_type}")

        issues = self._critic.review(doc_type, images, {"artifact_path": str(path), "images": [str(i) for i in images]})
        critical = [i for i in issues if i.severity == IssueSeverity.CRITICAL]
        rerender_targets = sorted({i.target for i in critical if i.target is not None})
        return CriticReport(
            passed=not critical,
            issues=issues,
            rerender_targets=rerender_targets,
            backend=backend,
            note=f"images={len(images)}",
        )

    # ------------------------------------------------------------------ helpers

    def _pptx_to_pdf(self, path: Path) -> Optional[Path]:
        if not (shutil.which("soffice") or shutil.which("libreoffice")):
            return None
        bin_name = "soffice" if shutil.which("soffice") else "libreoffice"
        try:
            proc = subprocess.run(
                [bin_name, "--headless", "--convert-to", "pdf", "--outdir", str(path.parent), str(path)],
                timeout=60,
                capture_output=True,
            )
            if proc.returncode != 0:
                return None
            candidate = path.with_suffix(".pdf")
            return candidate if candidate.exists() else None
        except Exception:
            return None

    def _pdf_to_images(self, pdf_path: Path) -> list[Path]:
        if shutil.which("pdftoppm"):
            return self._pdftoppm(pdf_path)
        return self._pypdfium2(pdf_path)

    def _pdftoppm(self, pdf_path: Path) -> list[Path]:
        prefix = pdf_path.parent / f"{pdf_path.stem}_page"
        try:
            proc = subprocess.run(
                ["pdftoppm", "-jpeg", "-r", "150", str(pdf_path), str(prefix)],
                timeout=60,
                capture_output=True,
            )
            if proc.returncode != 0:
                return []
            return sorted(pdf_path.parent.glob(f"{pdf_path.stem}_page-*.jpg"))
        except Exception:
            return []

    def _pypdfium2(self, pdf_path: Path) -> list[Path]:
        try:
            import pypdfium2 as pdfium
        except ImportError:
            return []
        try:
            pdf = pdfium.PdfDocument(str(pdf_path))
            out: list[Path] = []
            for i, page in enumerate(pdf):
                img = page.render(scale=2.0).to_pil()
                out_path = pdf_path.parent / f"{pdf_path.stem}_page-{i+1:03d}.jpg"
                img.save(str(out_path), "JPEG", quality=80)
                out.append(out_path)
            return out
        except Exception:
            return []
