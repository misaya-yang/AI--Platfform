"""PptxRenderer — design-system-driven.

Design conventions (distilled from 2026 SOTA research):

* Token-based colours (surface / ink / accent), not ``palette[0..4]``.
* Typography scale 1.333 (eyebrow / display / h1 / h2 / lead / body / caption).
* **Eyebrow kicker** above every H1 — uppercase, tracked, ~11pt, muted.
* **Section numeral** behind content — 180pt "01 / 08", ~20% opacity.
* **No accent underline under the title** — Anthropic's skill flags this
  as the #1 AI-tell.
* **Asymmetric splits** (38/62), not 50/50.
* **Soft drop shadow** on every elevated surface via ``a:effectLst`` XML.
* **Gradient surface tint** for hero and card backgrounds.
* **Background geometry** — rotated soft rectangle at 4-8% fill.
* **Layout variety** — caller's planner is responsible for no-two-same
  consecutive slides; renderer honours whatever IR ships.
"""

from __future__ import annotations

import re
import time
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ..design_system import (
    ColorTokens,
    DesignSystem,
    TypeScale,
    available_systems,
    design_system_from_palette,
    get_design_system,
    shade,
    text_on,
    tint,
)
from ..ir import (
    BulletBlock,
    ChartBlock,
    ChartSpec,
    HeadingBlock,
    ImageBlock,
    ParagraphBlock,
    PptxIR,
    PptxSlide,
    QuoteBlock,
    TableBlock,
)
from ..ir.pptx import _ImageSource, _IconRef, _ShapeSpec, VisualSpec  # noqa
from .base import BaseRenderer, RenderError, RenderResult
from .pptx_effects import (
    add_outer_shadow,
    set_linear_gradient_fill,
    set_no_stroke,
    set_rotation,
    set_stroke,
    set_transparency,
)


SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.66          # outer safe area (≈ 48pt)
GRID_GAP = 0.33


# ---------------------------------------------------------------------------
# colour helpers


def _rgb(hex_value: str) -> RGBColor:
    return RGBColor(
        int(hex_value[0:2], 16),
        int(hex_value[2:4], 16),
        int(hex_value[4:6], 16),
    )


_CHART_KINDS = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER_LINES,
}


# ---------------------------------------------------------------------------
# renderer


@dataclass
class _Ctx:
    """Per-render context — design system + typography + index."""

    ds: DesignSystem
    total_slides: int
    eyebrow_default: str
    brand: str
    author: str

    @property
    def c(self) -> ColorTokens:
        return self.ds.colors

    @property
    def t(self) -> TypeScale:
        return self.ds.type_scale


class PptxRenderer(BaseRenderer):
    format = "pptx"

    # ------------------------------------------------------------------ setup

    def _ctx(self, ir: PptxIR) -> _Ctx:
        # Pick design system: explicit metadata.extra["design_system"]
        # overrides; otherwise derive from the 5-colour palette bridge.
        ds_name = None
        if ir.metadata.model_extra:
            ds_name = ir.metadata.model_extra.get("design_system")
        ds: Optional[DesignSystem] = None
        if ds_name and ds_name in available_systems():
            ds = get_design_system(ds_name)
        if ds is None:
            palette_hex = [c.value for c in ir.theme.palette]
            ds = design_system_from_palette(palette_hex)

        subtitle = ir.metadata.subtitle or ""
        eyebrow = subtitle.upper() if subtitle else ""
        brand = ir.metadata.author or "PRESENTATION"
        return _Ctx(
            ds=ds,
            total_slides=len(ir.content.slides),
            eyebrow_default=eyebrow,
            brand=brand,
            author=ir.metadata.author or "",
        )

    def _new_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W)
        prs.slide_height = Inches(SLIDE_H)
        return prs

    # ----------------------------------------------------------- primitives

    def _rect(self, slide, *, left, top, width, height, hex_fill, shadow=False, radius=False):
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(hex_fill)
        set_no_stroke(shp)
        if radius:
            try:
                shp.adjustments[0] = 0.05
            except Exception:
                pass
        if shadow:
            add_outer_shadow(shp, blur_emu=40_000, dist_emu=23_000, alpha_per_mille=18_000)
        return shp

    def _gradient_rect(self, slide, *, left, top, width, height, stops, angle=90.0, shadow=False):
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        set_no_stroke(shp)
        set_linear_gradient_fill(shp, stops=stops, angle_deg=angle)
        if shadow:
            add_outer_shadow(shp, blur_emu=40_000, dist_emu=23_000, alpha_per_mille=18_000)
        return shp

    def _ellipse(self, slide, *, left, top, width, height, hex_fill, shadow=False):
        shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(width), Inches(height))
        shp.fill.solid()
        shp.fill.fore_color.rgb = _rgb(hex_fill)
        set_no_stroke(shp)
        if shadow:
            add_outer_shadow(shp, blur_emu=30_000, dist_emu=15_000, alpha_per_mille=15_000)
        return shp

    def _text(
        self,
        slide,
        text: str,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        size_pt: float,
        bold: bool = False,
        italic: bool = False,
        align: str = "left",
        v_anchor: str = "top",
        rgb_hex: str = "0F172A",
        font: str = "Calibri",
        line_spacing: Optional[float] = None,
        tracking_pct: float = 0.0,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.02)
        tf.margin_right = Inches(0.02)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(v_anchor, MSO_ANCHOR.TOP)
        p = tf.paragraphs[0]
        p.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }.get(align, PP_ALIGN.LEFT)
        if line_spacing is not None:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(rgb_hex)
        # Letter tracking via XML (a:rPr spc attribute, 100ths of point)
        if tracking_pct != 0.0:
            try:
                from lxml import etree
                _A = "http://schemas.openxmlformats.org/drawingml/2006/main"
                rpr = run._r.find(f"{{{_A}}}rPr")
                if rpr is not None:
                    spc_units = int(tracking_pct * 1000)  # ~em thousandths
                    rpr.set("spc", str(spc_units))
            except Exception:
                pass
        return box

    # ----------------------------------------------------------- chrome

    def _background(self, slide, hex_fill: str) -> None:
        bg = self._rect(slide, left=0, top=0, width=SLIDE_W, height=SLIDE_H, hex_fill=hex_fill)
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def _background_gradient(self, slide, from_hex: str, to_hex: str, angle: float = 135.0) -> None:
        bg = self._gradient_rect(slide, left=0, top=0, width=SLIDE_W, height=SLIDE_H, stops=[(0.0, from_hex), (1.0, to_hex)], angle=angle)
        spTree = bg._element.getparent()
        spTree.remove(bg._element)
        spTree.insert(2, bg._element)

    def _eyebrow(self, slide, text: str, *, left: float, top: float, ctx: _Ctx, width: float = 8.0):
        """Small uppercase, tracked text above an H1. The single biggest
        "looks designed" lever."""
        text = text.strip()
        if not text:
            return None
        # Small coloured square + label
        dot_size = 0.14
        self._rect(
            slide,
            left=left,
            top=top + 0.06,
            width=dot_size,
            height=dot_size,
            hex_fill=ctx.c.accent,
        )
        return self._text(
            slide,
            text.upper(),
            left=left + dot_size + 0.18,
            top=top,
            width=width,
            height=0.3,
            size_pt=ctx.t.eyebrow_pt,
            bold=True,
            rgb_hex=ctx.c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=ctx.t.eyebrow_tracking_pct,
        )

    def _section_numeral(self, slide, index: int, *, ctx: _Ctx, x: float, y: float) -> None:
        """Giant muted numeral in the corner — magazine-style decoration.

        Uses outline-ish appearance by drawing in ink_muted at low
        opacity. Size is ~180pt so it reads as pattern, not content.
        """
        label = f"{index:02d} / {ctx.total_slides:02d}"
        self._text(
            slide,
            label,
            left=x,
            top=y,
            width=4.5,
            height=2.5,
            size_pt=ctx.t.section_numeral_pt * 0.6,   # 108pt — looks "big" without overflowing
            bold=True,
            rgb_hex=tint(ctx.c.ink_muted, 0.6),
            font=ctx.ds.font_display,
            line_spacing=1.0,
            v_anchor="top",
        )

    def _footer(self, slide, *, ctx: _Ctx, index: int) -> None:
        """Small footer bar with brand + page number. No coloured strip —
        the trick is that it reads as metadata, not decoration."""
        # Left: brand
        self._text(
            slide,
            (ctx.brand or "").upper(),
            left=MARGIN,
            top=SLIDE_H - 0.42,
            width=SLIDE_W / 2,
            height=0.28,
            size_pt=9,
            bold=True,
            rgb_hex=ctx.c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=0.08,
        )
        # Right: page number
        self._text(
            slide,
            f"{index:02d} / {ctx.total_slides:02d}",
            left=SLIDE_W - MARGIN - 2.0,
            top=SLIDE_H - 0.42,
            width=2.0,
            height=0.28,
            size_pt=9,
            bold=True,
            align="right",
            rgb_hex=ctx.c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=0.08,
        )

    def _background_geometry(self, slide, *, ctx: _Ctx) -> None:
        """Subtle decoration — a rotated soft square and a corner quarter-circle.
        Both rendered at 5-8% alpha via transparency injection, so they read
        as faint pattern."""
        # Rotated faint square (bottom-left)
        sq = self._rect(
            slide,
            left=-1.5,
            top=SLIDE_H - 2.2,
            width=3.5,
            height=3.5,
            hex_fill=ctx.c.accent,
        )
        set_transparency(sq, alpha_per_mille=6_000)
        set_rotation(sq, degrees=15.0)

        # Corner circle (top-right)
        circ = self._ellipse(
            slide,
            left=SLIDE_W - 2.0,
            top=-2.0,
            width=4.0,
            height=4.0,
            hex_fill=ctx.c.accent_secondary,
        )
        set_transparency(circ, alpha_per_mille=5_000)

    # ----------------------------------------------------------- layouts

    def _draw_title(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        """Title slide — asymmetric split with coloured panel on left.

        The left 38% is a solid coloured band (surface_inverted) carrying
        brand/kicker/date. The right 62% is the light page surface
        carrying the big title — dark ink, easy to read. This design
        survives font substitution and QuickLook previews.
        """
        c = ctx.c
        t = ctx.t
        # Page surface
        self._background(slide, c.surface)

        # Left coloured panel
        panel_w = SLIDE_W * 0.38
        self._rect(slide, left=0, top=0, width=panel_w, height=SLIDE_H, hex_fill=c.surface_inverted)

        # Accent rule at top of panel
        self._rect(slide, left=0, top=0, width=panel_w, height=0.08, hex_fill=c.accent)

        # Eyebrow on left panel
        self._text(
            slide,
            (s.subtitle or ir.metadata.subtitle or "Keynote").upper(),
            left=0.55,
            top=0.55,
            width=panel_w - 1.0,
            height=0.35,
            size_pt=t.eyebrow_pt,
            bold=True,
            rgb_hex=c.accent,
            font=ctx.ds.font_body,
            tracking_pct=t.eyebrow_tracking_pct,
        )

        # Decorative number on left panel (big, muted) — wide enough for
        # two digits at 160pt in either serif or sans-serif.
        self._text(
            slide,
            "01",
            left=0.55,
            top=SLIDE_H - 4.3,
            width=panel_w - 1.0,
            height=3.2,
            size_pt=160,
            bold=True,
            rgb_hex=tint(c.surface_inverted, 0.08),
            font=ctx.ds.font_display,
            line_spacing=1.0,
        )

        # Brand bottom-left of panel
        self._text(
            slide,
            (ctx.brand or "PRESENTATION").upper(),
            left=0.55,
            top=SLIDE_H - 0.75,
            width=panel_w - 1.0,
            height=0.35,
            size_pt=10,
            bold=True,
            rgb_hex=tint(c.ink_inverted, 0.35),
            font=ctx.ds.font_body,
            tracking_pct=0.14,
        )

        # Right side: title on the light surface
        right_left = panel_w + 0.6
        right_w = SLIDE_W - right_left - MARGIN

        title_text = s.title or ir.metadata.title
        # Split the title on the last " — " if present for a 2-line look
        lines = title_text.split(" — ", 1) if " — " in title_text else [title_text]

        # Top kicker on right (more muted)
        self._text(
            slide,
            (ir.metadata.subtitle or "").upper() or "2026",
            left=right_left,
            top=0.7,
            width=right_w,
            height=0.35,
            size_pt=t.eyebrow_pt,
            bold=True,
            rgb_hex=c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=t.eyebrow_tracking_pct,
        )
        # Tiny accent dot
        self._rect(slide, left=right_left, top=1.15, width=0.32, height=0.08, hex_fill=c.accent)

        # Big title
        display_size = t.display_pt if len(title_text) < 40 else (t.display_pt - 12)
        self._text(
            slide,
            title_text,
            left=right_left,
            top=SLIDE_H / 2 - 1.8,
            width=right_w,
            height=4.0,
            size_pt=display_size,
            bold=True,
            rgb_hex=c.ink_primary,
            font=ctx.ds.font_display,
            line_spacing=1.08,
        )

        # Author/date + page number at bottom of right side
        meta_bits = [b for b in (ctx.author, ir.metadata.created_at) if b]
        meta_line = " · ".join(meta_bits)
        if meta_line:
            self._text(
                slide,
                meta_line,
                left=right_left,
                top=SLIDE_H - 0.9,
                width=right_w - 1.5,
                height=0.35,
                size_pt=t.body_pt - 2,
                rgb_hex=c.ink_muted,
                font=ctx.ds.font_body,
            )
        self._text(
            slide,
            f"{index:02d} / {ctx.total_slides:02d}",
            left=SLIDE_W - MARGIN - 1.5,
            top=SLIDE_H - 0.75,
            width=1.5,
            height=0.35,
            size_pt=10,
            bold=True,
            align="right",
            rgb_hex=c.ink_muted,
            font=ctx.ds.font_body,
        )

    def _draw_title_content(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface)

        # Eyebrow
        eyebrow_text = self._pick_eyebrow(s, ctx)
        self._eyebrow(slide, eyebrow_text, left=MARGIN, top=MARGIN, ctx=ctx)

        # H1
        title_text = s.title or ""
        self._text(
            slide,
            title_text,
            left=MARGIN,
            top=MARGIN + 0.45,
            width=SLIDE_W - 2 * MARGIN - 2.5,
            height=1.2,
            size_pt=t.h1_pt,
            bold=True,
            rgb_hex=c.ink_primary,
            font=ctx.ds.font_display,
            line_spacing=1.1,
        )

        # Big muted section numeral — top right
        self._text(
            slide,
            f"{index:02d}",
            left=SLIDE_W - MARGIN - 1.8,
            top=MARGIN - 0.2,
            width=1.8,
            height=1.8,
            size_pt=108,
            bold=True,
            align="right",
            rgb_hex=tint(c.ink_muted, 0.55),
            font=ctx.ds.font_display,
            line_spacing=1.0,
        )

        # Optional subtitle / lead paragraph
        cursor_top = MARGIN + 1.75
        if s.subtitle:
            self._text(
                slide,
                s.subtitle,
                left=MARGIN,
                top=cursor_top,
                width=SLIDE_W - 2 * MARGIN - 2.5,
                height=0.7,
                size_pt=t.lead_pt,
                rgb_hex=c.ink_secondary,
                font=ctx.ds.font_body,
                line_spacing=1.3,
            )
            cursor_top += 0.8

        body = s.body or []
        content_left = MARGIN
        content_width = SLIDE_W - 2 * MARGIN
        content_top = cursor_top + 0.25
        content_height = SLIDE_H - content_top - 0.9

        # Content-shape-driven layout upgrades. Priority order:
        #   1. single short paragraph (< 240 chars) → LEAD paragraph (big)
        #   2. single BulletBlock 2-6 items        → feature-row stack
        #   3. otherwise                           → normal block stack
        simple_bullets = (
            len(body) == 1
            and isinstance(body[0], BulletBlock)
            and 2 <= len(body[0].items) <= 6
        )
        single_short_para = (
            len(body) == 1
            and isinstance(body[0], ParagraphBlock)
            and len(body[0].text) <= 240
        )
        single_long_para = (
            len(body) == 1
            and isinstance(body[0], ParagraphBlock)
            and len(body[0].text) <= 600
        )

        if single_short_para:
            self._draw_lead_paragraph(slide, body[0].text, left=content_left, top=content_top, width=content_width, height=content_height, ctx=ctx)
        elif simple_bullets:
            self._draw_feature_rows(slide, body[0].items, ordered=body[0].ordered, left=content_left, top=content_top, width=content_width, height=content_height, ctx=ctx)
        elif single_long_para:
            self._draw_lead_paragraph(slide, body[0].text, left=content_left, top=content_top, width=content_width, height=content_height, ctx=ctx, size_pt=ctx.t.lead_pt + 4)
        else:
            self._draw_blocks(slide, body, left=content_left, top=content_top, width=content_width, height=content_height, ctx=ctx)

        self._footer(slide, ctx=ctx, index=index)

    def _draw_lead_paragraph(self, slide, text: str, *, left: float, top: float, width: float, height: float, ctx: _Ctx, size_pt: Optional[float] = None):
        """Render a single short paragraph as a HERO-sized lead statement.

        Uses editorial magazine "pull-quote" conventions:
          * Big left-aligned type (~28-44pt depending on length).
          * Left-edge accent bar (vertical line, palette.accent).
          * Decorative oversized initial character in palette.accent_subtle.
          * Caption row below with a thin muted rule, echoing the footer.

        Occupies the whole content region so the slide reads intentional.
        """
        c = ctx.c
        t = ctx.t
        n_chars = max(1, len(text))
        if size_pt is None:
            if n_chars < 80:
                size_pt = 44.0
            elif n_chars < 160:
                size_pt = 34.0
            elif n_chars < 240:
                size_pt = 28.0
            else:
                size_pt = 22.0

        # Big ornamental quote glyph (tinted, behind the text)
        self._text(
            slide,
            "\u201C",
            left=left,
            top=top - 0.4,
            width=1.8,
            height=2.4,
            size_pt=220,
            bold=True,
            rgb_hex=tint(c.accent, 0.72),
            font=ctx.ds.font_display,
            line_spacing=1.0,
        )

        # Vertical accent bar on the left edge of the text column
        text_indent = 1.2
        self._rect(
            slide,
            left=left + text_indent - 0.25,
            top=top + 0.4,
            width=0.06,
            height=min(height - 0.8, size_pt * 0.05 + 1.6),
            hex_fill=c.accent,
        )

        # Lead text (left of centre, line-length cap ≈ 60 chars)
        self._text(
            slide,
            text,
            left=left + text_indent,
            top=top + 0.35,
            width=min(width - text_indent, 10.0),
            height=max(3.0, size_pt * 0.07 + 2.5),
            size_pt=size_pt,
            bold=False,
            align="left",
            v_anchor="top",
            rgb_hex=c.ink_primary,
            font=ctx.ds.font_display,
            line_spacing=1.22,
        )

        # Thin muted divider near bottom of content area + caption
        caption_top = top + height - 0.75
        self._rect(
            slide,
            left=left + text_indent,
            top=caption_top,
            width=2.0,
            height=0.02,
            hex_fill=c.border_strong,
        )
        self._text(
            slide,
            "KEY TAKEAWAY",
            left=left + text_indent,
            top=caption_top + 0.12,
            width=5.0,
            height=0.35,
            size_pt=t.caption_pt - 1,
            bold=True,
            rgb_hex=c.ink_muted,
            font=ctx.ds.font_body,
            tracking_pct=0.14,
        )

    def _draw_two_col(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        """Asymmetric 38/62 split — hero text left, support column right."""
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface)
        self._eyebrow(slide, self._pick_eyebrow(s, ctx), left=MARGIN, top=MARGIN, ctx=ctx)

        # Title left column
        left_w = (SLIDE_W - 2 * MARGIN) * 0.38 - GRID_GAP / 2
        right_w = (SLIDE_W - 2 * MARGIN) * 0.62 - GRID_GAP / 2
        right_left = MARGIN + left_w + GRID_GAP

        self._text(
            slide,
            s.title or "",
            left=MARGIN,
            top=MARGIN + 0.45,
            width=left_w,
            height=3.5,
            size_pt=t.h1_pt,
            bold=True,
            rgb_hex=c.ink_primary,
            font=ctx.ds.font_display,
            line_spacing=1.1,
        )
        if s.subtitle:
            self._text(
                slide,
                s.subtitle,
                left=MARGIN,
                top=MARGIN + 0.45 + 2.0,
                width=left_w,
                height=2.0,
                size_pt=t.lead_pt,
                rgb_hex=c.ink_secondary,
                font=ctx.ds.font_body,
                line_spacing=1.3,
            )

        # Right card stack from body
        body = s.body or []
        card_h = (SLIDE_H - MARGIN - 0.9 - MARGIN - 0.45) / max(1, len(body))
        card_h = min(card_h, 1.6)
        cursor = MARGIN + 0.45
        for i, blk in enumerate(body[:4]):
            card = self._rect(
                slide,
                left=right_left,
                top=cursor,
                width=right_w,
                height=card_h - 0.18,
                hex_fill=c.surface_elevated,
                shadow=True,
                radius=True,
            )
            # top accent edge
            self._rect(
                slide,
                left=right_left,
                top=cursor,
                width=0.12,
                height=card_h - 0.18,
                hex_fill=c.accent if i == 0 else c.accent_secondary,
            )
            # content
            if isinstance(blk, ParagraphBlock):
                self._text(
                    slide, blk.text,
                    left=right_left + 0.35, top=cursor + 0.25,
                    width=right_w - 0.5, height=card_h - 0.68,
                    size_pt=t.body_pt, rgb_hex=c.ink_secondary,
                    font=ctx.ds.font_body, line_spacing=1.35,
                )
            elif isinstance(blk, BulletBlock):
                self._bullets(
                    slide, blk.items,
                    left=right_left + 0.35, top=cursor + 0.18,
                    width=right_w - 0.5, height=card_h - 0.48,
                    ctx=ctx, ordered=blk.ordered,
                )
            elif isinstance(blk, HeadingBlock):
                self._text(
                    slide, blk.text,
                    left=right_left + 0.35, top=cursor + 0.3,
                    width=right_w - 0.5, height=card_h - 0.68,
                    size_pt=t.h2_pt, bold=True, rgb_hex=c.ink_primary,
                    font=ctx.ds.font_display,
                )
            cursor += card_h

        self._footer(slide, ctx=ctx, index=index)

    def _draw_quote(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        """Dark full-bleed pull-quote slide.

        Composition:
          * Background: surface_inverted (deep).
          * Large circle bleeding in from the bottom-right corner,
            low-alpha accent — subtle, not competing with text.
          * Content column is centre-vertical: big "❝" glyph in accent,
            then the quote at 40pt, then an accent rule + author.
        """
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface_inverted)

        # Corner-bleed decorative circle (bottom-right). Use a shade-blended
        # accent so it sits naturally on the dark surface (LibreOffice drops
        # a:alpha transparency on solidFill, so we pre-mix instead).
        blended = shade(c.accent, 0.55)  # darken accent toward the dark bg
        self._ellipse(
            slide,
            left=SLIDE_W - 3.5, top=SLIDE_H - 3.5,
            width=6.0, height=6.0,
            hex_fill=blended,
        )

        # Thin top accent rule (brand marker)
        self._rect(slide, left=MARGIN, top=MARGIN, width=0.8, height=0.06, hex_fill=c.accent)

        # Extract quote + author from body
        quote_text = None
        author = None
        for b in s.body or []:
            if isinstance(b, QuoteBlock):
                quote_text = b.text
                author = b.author
                break
            if isinstance(b, ParagraphBlock) and quote_text is None:
                quote_text = b.text
        if s.title and not quote_text:
            quote_text = s.title

        # Big opening glyph — centered vertically above the text
        glyph_top = SLIDE_H / 2 - 2.6
        self._text(
            slide, "\u201C",
            left=MARGIN, top=glyph_top,
            width=2.4, height=2.4,
            size_pt=200, bold=True,
            rgb_hex=c.accent,
            font=ctx.ds.font_display,
            line_spacing=1.0,
        )

        # Quote body — centre vertically, left-aligned (no centred body)
        quote_top = glyph_top + 1.8
        text_w = SLIDE_W - 2 * MARGIN
        self._text(
            slide, quote_text or "",
            left=MARGIN, top=quote_top,
            width=min(text_w, 11.0), height=3.2,
            size_pt=40, bold=True, italic=False,
            rgb_hex=c.ink_inverted, font=ctx.ds.font_display,
            line_spacing=1.22,
        )

        # Accent rule + author line
        attr_top = quote_top + 3.2
        self._rect(slide, left=MARGIN, top=attr_top, width=0.7, height=0.06, hex_fill=c.accent)
        if author:
            self._text(
                slide, author.upper(),
                left=MARGIN + 0.9, top=attr_top - 0.1,
                width=10.0, height=0.4,
                size_pt=t.eyebrow_pt + 1, bold=True,
                rgb_hex=c.accent, font=ctx.ds.font_body,
                tracking_pct=0.14,
            )

        # Footer — brand + page
        self._text(
            slide, (ctx.brand or "").upper(),
            left=MARGIN, top=SLIDE_H - 0.5,
            width=8.0, height=0.3,
            size_pt=9, bold=True,
            rgb_hex=tint(c.ink_inverted, 0.45),
            font=ctx.ds.font_body,
            tracking_pct=0.12,
        )
        self._text(
            slide, f"{index:02d} / {ctx.total_slides:02d}",
            left=SLIDE_W - MARGIN - 1.5, top=SLIDE_H - 0.5,
            width=1.5, height=0.3, size_pt=9, align="right",
            bold=True, rgb_hex=tint(c.ink_inverted, 0.45),
            font=ctx.ds.font_body,
        )

    def _draw_stat_callout(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface)
        self._eyebrow(slide, self._pick_eyebrow(s, ctx), left=MARGIN, top=MARGIN, ctx=ctx)
        self._text(
            slide, s.title or "",
            left=MARGIN, top=MARGIN + 0.45,
            width=SLIDE_W - 2 * MARGIN, height=1.0,
            size_pt=t.h2_pt, bold=True,
            rgb_hex=c.ink_primary, font=ctx.ds.font_display,
        )

        # Big stat number
        stat = s.stat_value or ""
        self._text(
            slide, stat,
            left=MARGIN, top=MARGIN + 1.4,
            width=SLIDE_W - 2 * MARGIN, height=4.2,
            size_pt=t.stat_numeral_pt, bold=True,
            align="left", v_anchor="middle",
            rgb_hex=c.accent, font=ctx.ds.font_display,
            line_spacing=1.0,
        )
        # label
        if s.stat_label:
            # accent rule + label
            self._rect(slide, left=MARGIN, top=SLIDE_H - MARGIN - 1.1, width=0.6, height=0.05, hex_fill=c.accent)
            self._text(
                slide, s.stat_label,
                left=MARGIN + 0.8, top=SLIDE_H - MARGIN - 1.3,
                width=SLIDE_W - 2 * MARGIN - 0.8, height=1.2,
                size_pt=t.lead_pt, bold=True,
                rgb_hex=c.ink_secondary, font=ctx.ds.font_body,
                line_spacing=1.3,
            )
        self._footer(slide, ctx=ctx, index=index)

    def _draw_icon_row(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface)
        self._eyebrow(slide, self._pick_eyebrow(s, ctx), left=MARGIN, top=MARGIN, ctx=ctx)
        self._text(
            slide, s.title or "",
            left=MARGIN, top=MARGIN + 0.45,
            width=SLIDE_W - 2 * MARGIN, height=1.1,
            size_pt=t.h1_pt, bold=True,
            rgb_hex=c.ink_primary, font=ctx.ds.font_display,
            line_spacing=1.1,
        )

        items = self._collect_label_list(s)
        n = min(max(len(items), 1), 4)
        # Each item: circle + big word + short subtitle. Alternate circle fills.
        row_top = MARGIN + 2.2
        usable = SLIDE_W - 2 * MARGIN
        cell_w = usable / n
        circle_d = 1.5
        ink_tokens = [c.accent, c.accent_secondary, c.ink_primary, shade(c.accent, 0.2)]

        for i in range(n):
            cx = MARGIN + cell_w * i + (cell_w - circle_d) / 2
            bg = ink_tokens[i % 4]
            self._ellipse(slide, left=cx, top=row_top, width=circle_d, height=circle_d, hex_fill=bg, shadow=True)
            # initial letter
            letter = items[i][:1].upper() if i < len(items) and items[i] else str(i + 1)
            self._text(
                slide, letter,
                left=cx, top=row_top, width=circle_d, height=circle_d,
                size_pt=56, bold=True, align="center", v_anchor="middle",
                rgb_hex=text_on(bg), font=ctx.ds.font_display,
            )
            # label below — title + short description split
            label = items[i] if i < len(items) else ""
            split_re = re.compile(r"\s+(?:via|—|-|:)\s+")
            parts = split_re.split(label, maxsplit=1)
            head = parts[0].strip()
            tail = parts[1].strip() if len(parts) > 1 else ""
            self._text(
                slide, head,
                left=MARGIN + cell_w * i + 0.1, top=row_top + circle_d + 0.3,
                width=cell_w - 0.2, height=0.6,
                size_pt=t.h2_pt - 4, bold=True, align="center",
                rgb_hex=c.ink_primary, font=ctx.ds.font_display,
            )
            if tail:
                self._text(
                    slide, tail,
                    left=MARGIN + cell_w * i + 0.15, top=row_top + circle_d + 0.95,
                    width=cell_w - 0.3, height=1.1,
                    size_pt=t.caption_pt + 1, align="center",
                    rgb_hex=c.ink_muted, font=ctx.ds.font_body,
                    line_spacing=1.35,
                )
        self._footer(slide, ctx=ctx, index=index)

    def _draw_grid_2x2(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface)
        self._eyebrow(slide, self._pick_eyebrow(s, ctx), left=MARGIN, top=MARGIN, ctx=ctx)
        self._text(
            slide, s.title or "",
            left=MARGIN, top=MARGIN + 0.45,
            width=SLIDE_W - 2 * MARGIN, height=1.1,
            size_pt=t.h1_pt, bold=True,
            rgb_hex=c.ink_primary, font=ctx.ds.font_display,
        )

        items = self._collect_label_list(s)
        while len(items) < 4:
            items.append("")
        items = items[:4]

        grid_top = MARGIN + 1.9
        grid_h = SLIDE_H - grid_top - MARGIN - 0.3
        cell_w = (SLIDE_W - 2 * MARGIN - GRID_GAP) / 2
        cell_h = (grid_h - GRID_GAP) / 2
        coords = [
            (MARGIN, grid_top),
            (MARGIN + cell_w + GRID_GAP, grid_top),
            (MARGIN, grid_top + cell_h + GRID_GAP),
            (MARGIN + cell_w + GRID_GAP, grid_top + cell_h + GRID_GAP),
        ]
        edge_colors = [c.accent, c.accent_secondary, c.ink_primary, shade(c.accent, 0.2)]

        for i, (lx, ty) in enumerate(coords):
            card = self._rect(
                slide, left=lx, top=ty, width=cell_w, height=cell_h,
                hex_fill=c.surface_elevated, shadow=True, radius=True,
            )
            # top accent
            self._rect(slide, left=lx, top=ty, width=cell_w, height=0.1, hex_fill=edge_colors[i])

            # big index number at top
            self._text(
                slide, f"{i + 1:02d}",
                left=lx + 0.35, top=ty + 0.35,
                width=1.3, height=0.9,
                size_pt=36, bold=True,
                rgb_hex=edge_colors[i], font=ctx.ds.font_display,
            )

            # Split label
            label = items[i]
            split_re = re.compile(r"\s+(?:via|—|-|:)\s+")
            parts = split_re.split(label, maxsplit=1)
            head = parts[0].strip() if parts else label
            tail = parts[1].strip() if len(parts) > 1 else ""
            self._text(
                slide, head,
                left=lx + 0.35, top=ty + 1.25,
                width=cell_w - 0.7, height=0.9,
                size_pt=t.h2_pt - 2, bold=True,
                rgb_hex=c.ink_primary, font=ctx.ds.font_display,
                line_spacing=1.15,
            )
            if tail:
                self._text(
                    slide, tail,
                    left=lx + 0.35, top=ty + 2.2,
                    width=cell_w - 0.7, height=cell_h - 2.5,
                    size_pt=t.body_pt - 2,
                    rgb_hex=c.ink_muted, font=ctx.ds.font_body,
                    line_spacing=1.4,
                )
        self._footer(slide, ctx=ctx, index=index)

    def _draw_halfbleed_image(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface_elevated)
        half_w = SLIDE_W * 0.58

        # Left panel (text)
        self._rect(slide, left=0, top=0, width=half_w, height=SLIDE_H, hex_fill=c.surface)
        self._eyebrow(slide, self._pick_eyebrow(s, ctx), left=MARGIN, top=MARGIN, ctx=ctx)
        self._text(
            slide, s.title or "",
            left=MARGIN, top=MARGIN + 0.45,
            width=half_w - MARGIN - 0.5, height=2.4,
            size_pt=t.h1_pt, bold=True,
            rgb_hex=c.ink_primary, font=ctx.ds.font_display,
            line_spacing=1.12,
        )
        # accent rule under title
        self._rect(slide, left=MARGIN, top=MARGIN + 2.55, width=0.8, height=0.05, hex_fill=c.accent)

        self._draw_blocks(
            slide, s.body,
            left=MARGIN, top=MARGIN + 2.95,
            width=half_w - MARGIN - 0.5,
            height=SLIDE_H - MARGIN - 0.6 - (MARGIN + 2.95),
            ctx=ctx,
        )

        # Right panel: image or decorative stack
        placed = False
        if s.visual and isinstance(s.visual.source, _ImageSource) and s.visual.source.path:
            try:
                slide.shapes.add_picture(s.visual.source.path, Inches(half_w), Inches(0), height=Inches(SLIDE_H))
                placed = True
            except Exception:
                placed = False
        if not placed:
            # Gradient "hero" block on right
            self._gradient_rect(
                slide,
                left=half_w, top=0,
                width=SLIDE_W - half_w, height=SLIDE_H,
                stops=[(0.0, c.accent), (1.0, shade(c.accent, 0.3))],
                angle=135.0,
            )
            # Large decorative circle — use a darker accent for natural blend
            # (LibreOffice drops alpha on solidFill, so we pre-mix).
            self._ellipse(
                slide, left=half_w + 0.5, top=SLIDE_H / 2 - 2.3,
                width=4.6, height=4.6,
                hex_fill=tint(c.accent, 0.25),
            )
            # decorative big numeral
            self._text(
                slide, f"{index:02d}",
                left=half_w + 0.3, top=MARGIN,
                width=SLIDE_W - half_w - MARGIN, height=1.8,
                size_pt=84, bold=True, align="right",
                rgb_hex=tint(c.accent_on, 0.3),
                font=ctx.ds.font_display,
            )

        # Footer only on the left (text) side
        self._text(
            slide, (ctx.brand or "").upper(),
            left=MARGIN, top=SLIDE_H - 0.5,
            width=half_w - MARGIN, height=0.3,
            size_pt=9, bold=True,
            rgb_hex=c.ink_muted, font=ctx.ds.font_body,
            tracking_pct=0.08,
        )

    def _draw_chart(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        c = ctx.c
        t = ctx.t
        self._background(slide, c.surface)
        self._eyebrow(slide, self._pick_eyebrow(s, ctx), left=MARGIN, top=MARGIN, ctx=ctx)
        self._text(
            slide, s.title or "",
            left=MARGIN, top=MARGIN + 0.45,
            width=SLIDE_W - 2 * MARGIN - 2.5, height=1.1,
            size_pt=t.h1_pt, bold=True,
            rgb_hex=c.ink_primary, font=ctx.ds.font_display,
        )

        spec: Optional[ChartSpec] = None
        if s.visual and isinstance(s.visual.source, ChartSpec):
            spec = s.visual.source
        else:
            for b in s.body or []:
                if isinstance(b, ChartBlock):
                    spec = b.spec
                    break

        if spec is None:
            self._text(
                slide, "[chart data missing]",
                left=MARGIN, top=2.2,
                width=SLIDE_W - 2 * MARGIN, height=4.0,
                size_pt=20, rgb_hex=c.ink_muted, font=ctx.ds.font_body,
            )
            self._footer(slide, ctx=ctx, index=index)
            return

        # Elevated chart card
        card = self._rect(
            slide, left=MARGIN, top=MARGIN + 1.7,
            width=SLIDE_W - 2 * MARGIN, height=SLIDE_H - MARGIN - 1.7 - MARGIN - 0.4,
            hex_fill=c.surface_elevated,
            shadow=True, radius=True,
        )

        chart_data = CategoryChartData()
        chart_data.categories = spec.categories
        for ser in spec.series:
            chart_data.add_series(ser["name"], ser["values"])

        slide.shapes.add_chart(
            _CHART_KINDS[spec.chart_type],
            Inches(MARGIN + 0.3), Inches(MARGIN + 2.0),
            Inches(SLIDE_W - 2 * MARGIN - 0.6),
            Inches(SLIDE_H - MARGIN - 2.0 - MARGIN - 0.7),
            chart_data,
        )
        self._footer(slide, ctx=ctx, index=index)

    def _draw_blank(self, prs, slide, ir: PptxIR, s: PptxSlide, ctx: _Ctx, index: int):
        self._background(slide, ctx.c.surface)
        if s.body:
            self._draw_blocks(slide, s.body, left=MARGIN, top=MARGIN, width=SLIDE_W - 2 * MARGIN, height=SLIDE_H - 2 * MARGIN, ctx=ctx)

    # ----------------------------------------------------------- block render

    def _bullets(self, slide, items, *, left, top, width, height, ctx: _Ctx, ordered: bool = False):
        c = ctx.c
        t = ctx.t
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.05)
        tf.margin_right = Inches(0.05)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.4
            p.space_after = Pt(6)
            # marker
            marker = p.add_run()
            marker.text = f"{i + 1}.  " if ordered else "●  "
            marker.font.name = ctx.ds.font_body
            marker.font.size = Pt(t.body_pt)
            marker.font.bold = True
            marker.font.color.rgb = _rgb(c.accent)
            # content
            tr = p.add_run()
            tr.text = item
            tr.font.name = ctx.ds.font_body
            tr.font.size = Pt(t.body_pt)
            tr.font.color.rgb = _rgb(c.ink_secondary)

    def _draw_feature_rows(self, slide, items: list[str], *, ordered: bool, left: float, top: float, width: float, height: float, ctx: _Ctx):
        """Rich alternative to bullet list — each item is a numbered row.

        Layout per row (row_h tall, textboxes share vertical centering):

            ┌────┐
            │ 01 │  HEAD
            └────┘  tail description — small muted
                    ──── (thin accent rule, sits BETWEEN rows)
        """
        c = ctx.c
        t = ctx.t
        n = len(items)
        row_gap = 0.25                     # a bit of air between rows
        row_h = min(1.1, max(0.85, (height - row_gap * (n - 1)) / n))
        chip_w = 0.95
        chip_h = row_h - 0.15
        text_left = left + chip_w + 0.45
        text_w = width - chip_w - 0.55
        head_h = 0.5
        tail_h = max(0.0, row_h - head_h - 0.1)

        split_re = re.compile(r"\s+(?:via|—|-)\s+|[:：]\s*")
        chip_colors = [c.accent, c.accent_secondary, c.ink_primary, shade(c.accent, 0.2)]

        for i, raw in enumerate(items):
            ty = top + i * (row_h + row_gap)
            chip_color = chip_colors[i % len(chip_colors)]
            self._rect(
                slide,
                left=left, top=ty + 0.04,
                width=chip_w, height=chip_h,
                hex_fill=chip_color,
                shadow=True, radius=True,
            )
            self._text(
                slide,
                f"{i + 1:02d}" if ordered else f"{i + 1:02d}",
                left=left, top=ty + 0.04,
                width=chip_w, height=chip_h,
                size_pt=26, bold=True,
                align="center", v_anchor="middle",
                rgb_hex=text_on(chip_color),
                font=ctx.ds.font_display,
            )
            parts = split_re.split(raw, maxsplit=1)
            head = parts[0].strip()
            tail = parts[1].strip() if len(parts) > 1 else ""
            self._text(
                slide, head,
                left=text_left, top=ty + 0.02,
                width=text_w, height=head_h,
                size_pt=t.h2_pt - 2, bold=True,
                rgb_hex=c.ink_primary, font=ctx.ds.font_display,
                line_spacing=1.1,
            )
            if tail:
                self._text(
                    slide, tail,
                    left=text_left, top=ty + 0.02 + head_h,
                    width=text_w, height=tail_h,
                    size_pt=t.body_pt - 1,
                    rgb_hex=c.ink_muted, font=ctx.ds.font_body,
                    line_spacing=1.35,
                )
            # Divider rule — sits in the GAP BELOW each row, never on text.
            # Skip the last row so we don't draw a trailing rule.
            if i < n - 1:
                self._rect(
                    slide,
                    left=text_left,
                    top=ty + row_h + row_gap / 2 - 0.01,
                    width=width - chip_w - 0.55,
                    height=0.02,
                    hex_fill=c.border_subtle,
                )

    def _draw_blocks(self, slide, blocks, *, left, top, width, height, ctx: _Ctx):
        c = ctx.c
        t = ctx.t
        if not blocks:
            return
        cursor = top
        remaining = height
        for b in blocks:
            if remaining <= 0.25:
                break
            h = 0.0
            if isinstance(b, HeadingBlock):
                h = 0.55
                self._text(
                    slide, b.text,
                    left=left, top=cursor,
                    width=width, height=h,
                    size_pt=t.h2_pt - 2, bold=True,
                    rgb_hex=c.ink_primary, font=ctx.ds.font_display,
                )
            elif isinstance(b, ParagraphBlock):
                n_chars = max(1, len(b.text))
                chars_per_line = max(60, int(width * 9))
                lines = max(1, (n_chars + chars_per_line - 1) // chars_per_line)
                h = min(remaining, 0.3 + lines * 0.36)
                self._text(
                    slide, b.text,
                    left=left, top=cursor,
                    width=width, height=h,
                    size_pt=t.body_pt,
                    align=b.align,
                    rgb_hex=c.ink_secondary,
                    font=ctx.ds.font_body,
                    line_spacing=1.4,
                )
            elif isinstance(b, BulletBlock):
                h = min(remaining, 0.3 + 0.4 * len(b.items))
                self._bullets(
                    slide, b.items,
                    left=left, top=cursor,
                    width=width, height=h,
                    ctx=ctx, ordered=b.ordered,
                )
            elif isinstance(b, QuoteBlock):
                h = min(remaining, 1.2)
                self._rect(slide, left=left, top=cursor + 0.1, width=0.08, height=h - 0.2, hex_fill=c.accent)
                self._text(
                    slide, f"\u201C{b.text}\u201D",
                    left=left + 0.25, top=cursor,
                    width=width - 0.25, height=h,
                    size_pt=t.lead_pt, italic=True,
                    rgb_hex=c.ink_primary, font=ctx.ds.font_display,
                    line_spacing=1.3,
                )
            elif isinstance(b, TableBlock):
                h = min(remaining, 0.45 * len(b.rows) + 0.3)
                self._draw_table(slide, b, left=left, top=cursor, width=width, height=h, ctx=ctx)
            elif isinstance(b, ImageBlock):
                if b.source_path and Path(b.source_path).exists():
                    h = min(remaining, 3.0)
                    try:
                        slide.shapes.add_picture(b.source_path, Inches(left), Inches(cursor), height=Inches(h))
                    except Exception:
                        h = 0.4
                        self._text(slide, f"[image: {b.alt_text}]", left=left, top=cursor, width=width, height=h, size_pt=13, italic=True, rgb_hex=c.ink_muted, font=ctx.ds.font_body)
                else:
                    h = 0.45
                    self._text(slide, f"[image: {b.alt_text}]", left=left, top=cursor, width=width, height=h, size_pt=13, italic=True, rgb_hex=c.ink_muted, font=ctx.ds.font_body)
            elif isinstance(b, ChartBlock):
                h = min(remaining, 3.2)
                try:
                    chart_data = CategoryChartData()
                    chart_data.categories = b.spec.categories
                    for ser in b.spec.series:
                        chart_data.add_series(ser["name"], ser["values"])
                    slide.shapes.add_chart(_CHART_KINDS[b.spec.chart_type], Inches(left), Inches(cursor), Inches(width), Inches(h), chart_data)
                except Exception:
                    self._text(slide, f"[chart: {b.alt_text}]", left=left, top=cursor, width=width, height=0.4, size_pt=13, rgb_hex=c.ink_muted, font=ctx.ds.font_body)
            else:
                h = 0.4
                self._text(slide, f"[{type(b).__name__}]", left=left, top=cursor, width=width, height=h, size_pt=12, rgb_hex=c.ink_muted, font=ctx.ds.font_body)
            cursor += h + 0.18
            remaining -= h + 0.18

    def _draw_table(self, slide, block: TableBlock, *, left, top, width, height, ctx: _Ctx):
        c = ctx.c
        rows = len(block.rows)
        cols = max(len(r.cells) for r in block.rows)
        shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
        table = shape.table
        for r_idx, row in enumerate(block.rows):
            for c_idx in range(cols):
                cell = table.cell(r_idx, c_idx)
                if c_idx < len(row.cells):
                    src = row.cells[c_idx]
                    cell.text = src.text
                    for p in cell.text_frame.paragraphs:
                        for run in p.runs:
                            run.font.size = Pt(11)
                            run.font.bold = src.bold or row.is_header
                            run.font.name = ctx.ds.font_body
                            if row.is_header:
                                run.font.color.rgb = _rgb(c.ink_inverted)
                            else:
                                run.font.color.rgb = _rgb(c.ink_secondary)
                    if row.is_header:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(c.ink_primary)
                    elif r_idx % 2 == 0:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = _rgb(c.surface_elevated)

    # ----------------------------------------------------------- helpers

    def _collect_label_list(self, s: PptxSlide) -> list[str]:
        items: list[str] = []
        for b in s.body or []:
            if isinstance(b, BulletBlock):
                items.extend(b.items)
            elif isinstance(b, ParagraphBlock):
                items.append(b.text)
            elif isinstance(b, HeadingBlock):
                items.append(b.text)
        return items

    def _pick_eyebrow(self, s: PptxSlide, ctx: _Ctx) -> str:
        # Planner may stash eyebrow in notes like "EYEBROW: ...".
        if s.notes and "EYEBROW:" in s.notes:
            for line in s.notes.splitlines():
                if line.strip().upper().startswith("EYEBROW:"):
                    return line.split(":", 1)[1].strip()
        if s.subtitle and len(s.subtitle) < 40:
            return s.subtitle
        if ctx.eyebrow_default:
            return ctx.eyebrow_default
        return "CHAPTER"

    # ----------------------------------------------------------- main

    async def render(self, ir: PptxIR, out_dir: Path) -> RenderResult:
        if not isinstance(ir, PptxIR):
            raise RenderError(f"PptxRenderer got wrong IR type: {type(ir).__name__}")
        started = time.perf_counter()
        prs = self._new_presentation()
        ctx = self._ctx(ir)

        dispatch = {
            "title": self._draw_title,
            "title_content": self._draw_title_content,
            "two_col": self._draw_two_col,
            "quote": self._draw_quote,
            "stat_callout": self._draw_stat_callout,
            "icon_row": self._draw_icon_row,
            "grid_2x2": self._draw_grid_2x2,
            "halfbleed_image": self._draw_halfbleed_image,
            "chart": self._draw_chart,
            "blank": self._draw_blank,
        }

        for idx, s_ir in enumerate(ir.content.slides, start=1):
            blank = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
            slide = prs.slides.add_slide(blank)
            fn = dispatch.get(s_ir.layout, self._draw_title_content)
            fn(prs, slide, ir, s_ir, ctx, idx)
            if s_ir.notes:
                slide.notes_slide.notes_text_frame.text = s_ir.notes

        out_dir.mkdir(parents=True, exist_ok=True)
        safe_name = re.sub(r"[^A-Za-z0-9_\-\.]", "_", ir.metadata.title)[:80] or "presentation"
        path = out_dir / f"{safe_name}.pptx"
        prs.save(str(path))

        return RenderResult(
            path=path,
            doc_type="pptx",
            bytes_size=path.stat().st_size,
            duration_ms=int((time.perf_counter() - started) * 1000),
            extra={"slides": len(ir.content.slides), "design_system": ctx.ds.name},
        )

    async def fix(self, ir: PptxIR, critic_findings, out_dir: Path) -> RenderResult:
        return await self.render(ir, out_dir)
