"""Block-level renderers used by multiple slide layouts.

These are the atoms that sit *inside* a layout's body region — bullets,
feature rows, lead paragraph, body block stack (light / dark), tables.
They take a ``prims`` Primitives instance and a bound rectangle.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path

from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from ...design_system import shade, text_on, tint
from ...ir import (
    BulletBlock,
    ChartBlock,
    HeadingBlock,
    ImageBlock,
    ParagraphBlock,
    QuoteBlock,
    TableBlock,
)
from ..primitives import rgb

logger = logging.getLogger(__name__)


CHART_KINDS = {
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER_LINES,
}


def bullets(slide, items, *, left, top, width, height, ctx, ordered=False):
    c = ctx.c
    t = ctx.t
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.4
        p.space_after = Pt(6)
        marker = p.add_run()
        marker.text = f"{i + 1}.  " if ordered else "●  "
        marker.font.name = ctx.ds.font_body
        marker.font.size = Pt(t.body_pt)
        marker.font.bold = True
        marker.font.color.rgb = rgb(c.accent)
        tr = p.add_run()
        tr.text = item
        tr.font.name = ctx.ds.font_body
        tr.font.size = Pt(t.body_pt)
        tr.font.color.rgb = rgb(c.ink_secondary)


def draw_feature_rows(slide, items, *, ordered, left, top, width, height, ctx, prims):
    """Numbered-chip row list. Each row: chip + head + tail + divider."""
    c = ctx.c
    t = ctx.t
    n = len(items)
    row_gap = 0.25
    row_h = min(1.1, max(0.85, (height - row_gap * (n - 1)) / n))
    chip_w = 0.95
    chip_h = row_h - 0.15
    text_left = left + chip_w + 0.45
    text_w = width - chip_w - 0.55
    head_h = 0.5
    tail_h = max(0.0, row_h - head_h - 0.1)

    split_re = re.compile(r"\s+(?:via|—|-)\s+|[:：]\s*")
    chip_colors = [c.accent, c.accent_secondary, c.ink_primary, shade(c.accent, 0.2)]

    for i, raw in enumerate(items):
        ty = top + i * (row_h + row_gap)
        chip_color = chip_colors[i % len(chip_colors)]
        prims.rect(slide, left=left, top=ty + 0.04, width=chip_w, height=chip_h,
                   hex_fill=chip_color, shadow=True, radius=True)
        prims.text(
            slide, f"{i + 1:02d}",
            left=left, top=ty + 0.04, width=chip_w, height=chip_h,
            size_pt=26, bold=True, align="center", v_anchor="middle",
            rgb_hex=text_on(chip_color), font=ctx.ds.font_display,
        )
        parts = split_re.split(raw, maxsplit=1)
        head = parts[0].strip()
        tail = parts[1].strip() if len(parts) > 1 else ""
        prims.text(
            slide, head,
            left=text_left, top=ty + 0.02, width=text_w, height=head_h,
            size_pt=t.h2_pt - 2, bold=True,
            rgb_hex=c.ink_primary, font=ctx.ds.font_display,
            line_spacing=1.1,
        )
        if tail:
            prims.text(
                slide, tail,
                left=text_left, top=ty + 0.02 + head_h, width=text_w, height=tail_h,
                size_pt=t.body_pt - 1,
                rgb_hex=c.ink_muted, font=ctx.ds.font_body,
                line_spacing=1.35,
            )
        if i < n - 1:
            prims.rect(
                slide,
                left=text_left, top=ty + row_h + row_gap / 2 - 0.01,
                width=width - chip_w - 0.55, height=0.02,
                hex_fill=c.border_subtle,
            )


def draw_lead_paragraph(slide, text, *, left, top, width, height, ctx, prims, size_pt=None):
    """Editorial pull-quote-sized single paragraph."""
    c = ctx.c
    t = ctx.t
    n_chars = max(1, len(text))
    if size_pt is None:
        if n_chars < 80:
            size_pt = 44.0
        elif n_chars < 160:
            size_pt = 34.0
        elif n_chars < 240:
            size_pt = 28.0
        else:
            size_pt = 22.0

    prims.text(
        slide, "“",
        left=left, top=top - 0.4, width=1.8, height=2.4,
        size_pt=220, bold=True,
        rgb_hex=tint(c.accent, 0.72),
        font=ctx.ds.font_display, line_spacing=1.0,
    )
    text_indent = 1.2
    prims.rect(
        slide,
        left=left + text_indent - 0.25, top=top + 0.4,
        width=0.06, height=min(height - 0.8, size_pt * 0.05 + 1.6),
        hex_fill=c.accent,
    )
    prims.text(
        slide, text,
        left=left + text_indent, top=top + 0.35,
        width=min(width - text_indent, 10.0),
        height=max(3.0, size_pt * 0.07 + 2.5),
        size_pt=size_pt, bold=False,
        align="left", v_anchor="top",
        rgb_hex=c.ink_primary, font=ctx.ds.font_display,
        line_spacing=1.22,
    )
    caption_top = top + height - 0.75
    prims.rect(
        slide,
        left=left + text_indent, top=caption_top,
        width=2.0, height=0.02,
        hex_fill=c.border_strong,
    )
    prims.text(
        slide, "KEY TAKEAWAY",
        left=left + text_indent, top=caption_top + 0.12,
        width=5.0, height=0.35,
        size_pt=t.caption_pt - 1, bold=True,
        rgb_hex=c.ink_muted, font=ctx.ds.font_body,
        tracking_pct=0.14,
    )


def draw_table(slide, block, *, left, top, width, height, ctx):
    c = ctx.c
    rows = len(block.rows)
    cols = max(len(r.cells) for r in block.rows)
    shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    for r_idx, row in enumerate(block.rows):
        for c_idx in range(cols):
            cell = table.cell(r_idx, c_idx)
            if c_idx < len(row.cells):
                src = row.cells[c_idx]
                cell.text = src.text
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(11)
                        run.font.bold = src.bold or row.is_header
                        run.font.name = ctx.ds.font_body
                        if row.is_header:
                            run.font.color.rgb = rgb(c.ink_inverted)
                        else:
                            run.font.color.rgb = rgb(c.ink_secondary)
                if row.is_header:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(c.ink_primary)
                elif r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(c.surface_elevated)


def draw_blocks(slide, blocks, *, left, top, width, height, ctx, prims):
    c = ctx.c
    t = ctx.t
    if not blocks:
        return
    cursor = top
    remaining = height
    for b in blocks:
        if remaining <= 0.25:
            break
        h = 0.0
        if isinstance(b, HeadingBlock):
            h = 0.55
            prims.text(
                slide, b.text,
                left=left, top=cursor, width=width, height=h,
                size_pt=t.h2_pt - 2, bold=True,
                rgb_hex=c.ink_primary, font=ctx.ds.font_display,
            )
        elif isinstance(b, ParagraphBlock):
            n_chars = max(1, len(b.text))
            chars_per_line = max(60, int(width * 9))
            lines = max(1, (n_chars + chars_per_line - 1) // chars_per_line)
            h = min(remaining, 0.3 + lines * 0.36)
            prims.text(
                slide, b.text,
                left=left, top=cursor, width=width, height=h,
                size_pt=t.body_pt, align=b.align,
                rgb_hex=c.ink_secondary, font=ctx.ds.font_body,
                line_spacing=1.4,
            )
        elif isinstance(b, BulletBlock):
            h = min(remaining, 0.3 + 0.4 * len(b.items))
            bullets(slide, b.items, left=left, top=cursor, width=width, height=h,
                    ctx=ctx, ordered=b.ordered)
        elif isinstance(b, QuoteBlock):
            h = min(remaining, 1.2)
            prims.rect(slide, left=left, top=cursor + 0.1, width=0.08, height=h - 0.2, hex_fill=c.accent)
            prims.text(
                slide, f"“{b.text}”",
                left=left + 0.25, top=cursor, width=width - 0.25, height=h,
                size_pt=t.lead_pt, italic=True,
                rgb_hex=c.ink_primary, font=ctx.ds.font_display,
                line_spacing=1.3,
            )
        elif isinstance(b, TableBlock):
            h = min(remaining, 0.45 * len(b.rows) + 0.3)
            draw_table(slide, b, left=left, top=cursor, width=width, height=h, ctx=ctx)
        elif isinstance(b, ImageBlock):
            if b.source_path and Path(b.source_path).exists():
                h = min(remaining, 3.0)
                try:
                    slide.shapes.add_picture(b.source_path, Inches(left), Inches(cursor), height=Inches(h))
                except (FileNotFoundError, OSError, ValueError) as exc:
                    logger.warning("block image %r failed to embed: %s", b.source_path, exc)
                    h = 0.4
                    prims.text(slide, f"[image: {b.alt_text}]", left=left, top=cursor,
                               width=width, height=h, size_pt=13, italic=True,
                               rgb_hex=c.ink_muted, font=ctx.ds.font_body)
            else:
                h = 0.45
                prims.text(slide, f"[image: {b.alt_text}]", left=left, top=cursor,
                           width=width, height=h, size_pt=13, italic=True,
                           rgb_hex=c.ink_muted, font=ctx.ds.font_body)
        elif isinstance(b, ChartBlock):
            h = min(remaining, 3.2)
            try:
                chart_data = CategoryChartData()
                chart_data.categories = b.spec.categories
                for ser in b.spec.series:
                    chart_data.add_series(ser["name"], ser["values"])
                slide.shapes.add_chart(CHART_KINDS[b.spec.chart_type],
                                      Inches(left), Inches(cursor), Inches(width), Inches(h),
                                      chart_data)
            except (KeyError, ValueError, TypeError) as exc:
                logger.warning("chart block render failed: %s", exc)
                prims.text(slide, f"[chart: {b.alt_text}]", left=left, top=cursor,
                           width=width, height=0.4, size_pt=13,
                           rgb_hex=c.ink_muted, font=ctx.ds.font_body)
        else:
            h = 0.4
            prims.text(slide, f"[{type(b).__name__}]", left=left, top=cursor,
                       width=width, height=h, size_pt=12,
                       rgb_hex=c.ink_muted, font=ctx.ds.font_body)
        cursor += h + 0.18
        remaining -= h + 0.18


def draw_blocks_dark(slide, blocks, *, left, top, width, height, ctx, prims):
    """Inverted-colour variant of ``draw_blocks``."""
    c = ctx.c
    t = ctx.t
    cursor = top
    remaining = height
    for b in blocks:
        if remaining <= 0.25:
            break
        h = 0.0
        if isinstance(b, HeadingBlock):
            h = 0.55
            prims.text(
                slide, b.text,
                left=left, top=cursor, width=width, height=h,
                size_pt=t.h2_pt - 2, bold=True,
                rgb_hex=c.ink_inverted, font=ctx.ds.font_display,
            )
        elif isinstance(b, ParagraphBlock):
            n_chars = max(1, len(b.text))
            chars_per_line = max(60, int(width * 9))
            lines = max(1, (n_chars + chars_per_line - 1) // chars_per_line)
            h = min(remaining, 0.3 + lines * 0.36)
            prims.text(
                slide, b.text,
                left=left, top=cursor, width=width, height=h,
                size_pt=t.body_pt,
                rgb_hex=tint(c.ink_inverted, 0.15),
                font=ctx.ds.font_body, line_spacing=1.4,
            )
        elif isinstance(b, BulletBlock):
            h = min(remaining, 0.3 + 0.4 * len(b.items))
            box = slide.shapes.add_textbox(Inches(left), Inches(cursor), Inches(width), Inches(h))
            tf = box.text_frame
            tf.word_wrap = True
            try:
                tf.auto_size = MSO_AUTO_SIZE.NONE
            except (AttributeError, ValueError) as exc:
                logger.debug("auto_size=NONE unavailable on dark bullet: %s", exc)
            for i, item in enumerate(b.items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.line_spacing = 1.4
                p.space_after = Pt(6)
                m = p.add_run()
                m.text = f"{i + 1}.  " if b.ordered else "●  "
                m.font.name = ctx.ds.font_body
                m.font.size = Pt(t.body_pt)
                m.font.bold = True
                m.font.color.rgb = rgb(c.accent)
                r = p.add_run()
                r.text = item
                r.font.name = ctx.ds.font_body
                r.font.size = Pt(t.body_pt)
                r.font.color.rgb = rgb(tint(c.ink_inverted, 0.15))
        elif isinstance(b, QuoteBlock):
            h = min(remaining, 1.2)
            prims.text(
                slide, f"“{b.text}”",
                left=left + 0.25, top=cursor, width=width - 0.25, height=h,
                size_pt=t.lead_pt, italic=True,
                rgb_hex=c.ink_inverted, font=ctx.ds.font_display,
                line_spacing=1.3,
            )
        else:
            h = 0.4
        cursor += h + 0.18
        remaining -= h + 0.18
