"""Layout-bounds regression tests for PptxRenderer.

These tests encode invariants that STRUCTURAL tests (which just check
"file exists and validates") miss. They would have caught the
`grid_2x2` cell-overflow bug and the title-truncation bug.

Three invariants checked:

1. **Shape-within-slide**: Every shape's ``top+height`` and ``left+width``
   sit inside the slide rectangle. A shape going negative or off the
   right/bottom indicates a layout math mistake.

2. **Shape-within-card (grid_2x2)**: In ``grid_2x2``, every shape whose
   top-left lies inside a card's bounds must also have its bottom-right
   inside that card (no leakage into footer or neighbouring cards).

3. **Stress: long content does not crash**: Rendering each layout with
   deliberately long CJK/ASCII content must not raise.

These are fast (no LibreOffice needed) and catch the specific family of
bugs that burned us.
"""

from __future__ import annotations

import asyncio
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from src.services.assistant.docgen.ir import (
    BulletBlock,
    ParagraphBlock,
    PptxContent,
    PptxIR,
    PptxSlide,
    Theme,
)
from src.services.assistant.docgen.ir.base import DocMetadata
from src.services.assistant.docgen.renderers.pptx_renderer import PptxRenderer


SLIDE_W_EMU = int(13.333 * 914400)
SLIDE_H_EMU = int(7.5 * 914400)

# Content that stresses every layout's geometry
LONG_PARA = (
    "输入（模式、PII、阻止列表、10ms 内的越狱检测），过程（置信度门控、范围执行、步骤验证），"
    "输出（格式校验、引用核实、禁用短语过滤）。护栏属于约束框架而非提示词——可独立更新部署。"
)
LONG_BULLETS = [
    "输入验证：模式检查、PII 检测、阻止列表、10ms 内越狱检测",
    "过程控制：置信度门控、范围执行、步骤验证",
    "输出验证：格式校验、引用核实、禁用短语过滤",
    "Sources: Anthropic safety research, Parallel.ai guardrail patterns",
]


def _make_ir(layout: str, *, body: list) -> PptxIR:
    return PptxIR(
        metadata=DocMetadata(title="Layout Stress Test", locale="en-US"),
        theme=Theme(),
        content=PptxContent(slides=[
            PptxSlide(layout="title", title="Stress Test"),
            PptxSlide(layout=layout, title="Layer 3: Guardrails", body=body),
        ]),
    )


async def _render(ir: PptxIR, tmp_path: Path) -> Path:
    return (await PptxRenderer().render(ir, tmp_path)).path


def _shapes_with_bounds(pptx_path: Path):
    """Yield (slide_idx, shape, (left, top, width, height)) for every
    shape — in Inches, as ints for pixel-accurate comparison."""
    prs = Presentation(str(pptx_path))
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and not shape.text_frame.text.strip():
                # Skip empty text frames — they're usually layout scaffolding
                continue
            yield si, shape, (shape.left, shape.top, shape.width, shape.height)


# -------------------------------------------------------------- invariant 1

@pytest.mark.asyncio
async def test_all_shapes_inside_slide(tmp_path: Path) -> None:
    """Every shape on every slide must be within the page rectangle."""
    ir = _make_ir("grid_2x2", body=[BulletBlock(items=LONG_BULLETS)])
    path = await _render(ir, tmp_path)
    violations = []
    for si, shape, (lx, ty, w, h) in _shapes_with_bounds(path):
        if lx < 0 or ty < 0:
            violations.append(f"slide {si+1}: negative origin ({lx}, {ty})")
        if lx + w > SLIDE_W_EMU + 10000:  # 10k EMU tolerance (~0.01")
            violations.append(f"slide {si+1}: shape overflows right edge ({lx + w} > {SLIDE_W_EMU})")
        if ty + h > SLIDE_H_EMU + 10000:
            violations.append(f"slide {si+1}: shape overflows bottom edge ({ty + h} > {SLIDE_H_EMU})")
    assert not violations, "Shape bounds violations:\n  " + "\n  ".join(violations)


# -------------------------------------------------------------- invariant 2

@pytest.mark.asyncio
async def test_grid_2x2_cell_containment(tmp_path: Path) -> None:
    """In grid_2x2, any textbox inside a card's left-half must have its
    bottom-right also inside that card's bounds (no leak into footer or
    next card).

    The grid cells are (in Inches, matching renderer constants):
      cell_w = (13.333 - 2*0.66 - 0.33) / 2 = 5.84
      grid_top = 2.56; grid_h = 3.98
      cell_h = (3.98 - 0.33) / 2 = 1.825

    Four cells at:
      TL (0.66, 2.56), TR (6.83, 2.56)
      BL (0.66, 4.72), BR (6.83, 4.72)   (all + cell_w × cell_h)
    """
    ir = _make_ir("grid_2x2", body=[BulletBlock(items=LONG_BULLETS)])
    path = await _render(ir, tmp_path)

    cells = [
        (Inches(0.66), Inches(2.56), Inches(5.84), Inches(1.825)),
        (Inches(0.66 + 5.84 + 0.33), Inches(2.56), Inches(5.84), Inches(1.825)),
        (Inches(0.66), Inches(2.56 + 1.825 + 0.33), Inches(5.84), Inches(1.825)),
        (Inches(0.66 + 5.84 + 0.33), Inches(2.56 + 1.825 + 0.33), Inches(5.84), Inches(1.825)),
    ]

    violations = []
    for si, shape, (lx, ty, w, h) in _shapes_with_bounds(path):
        if si != 1:  # only the grid_2x2 slide
            continue
        # Is this shape's top-left inside any cell?
        for idx, (cx, cy, cw, ch) in enumerate(cells):
            if cx <= lx < cx + cw and cy <= ty < cy + ch:
                # Bottom-right must also fit (20 000 EMU tolerance for shadow/radius)
                if lx + w > cx + cw + 20_000 or ty + h > cy + ch + 20_000:
                    violations.append(
                        f"cell {idx+1}: shape at ({lx},{ty}) size ({w},{h}) "
                        f"overflows cell ({cx+cw},{cy+ch})"
                    )
                break
    assert not violations, (
        "grid_2x2 cell containment violations:\n  " + "\n  ".join(violations)
    )


# -------------------------------------------------------------- invariant 3

@pytest.mark.asyncio
@pytest.mark.parametrize("layout", [
    "title_content", "two_col", "grid_2x2", "icon_row", "stat_callout",
    "quote", "halfbleed_image", "section_divider", "big_stat",
    "comparison_table", "pull_quote_card", "dark_card", "numbered_flow",
])
async def test_layout_does_not_crash_on_long_content(layout: str, tmp_path: Path) -> None:
    """Every layout must render without raising when given oversized content."""
    body = [
        ParagraphBlock(text=LONG_PARA),
        BulletBlock(items=LONG_BULLETS),
    ]
    ir = _make_ir(layout, body=body)
    # Some layouts need stat fields
    if layout in ("stat_callout", "big_stat"):
        ir.content.slides[1].stat_value = "1300"
        ir.content.slides[1].stat_label = "AI PRs / week"
    path = await _render(ir, tmp_path)
    assert path.exists()
    assert path.stat().st_size > 5_000, f"{layout}: pptx suspiciously small"


# -------------------------------------------------------------- invariant 4
# Bonus: auto_size=NONE actually applied

@pytest.mark.asyncio
async def test_textboxes_have_auto_size_none(tmp_path: Path) -> None:
    """Every textbox (except tables / placeholders) must have auto_size=NONE
    so text cannot silently grow its box and overflow."""
    from pptx.enum.text import MSO_AUTO_SIZE

    ir = _make_ir("grid_2x2", body=[BulletBlock(items=LONG_BULLETS)])
    path = await _render(ir, tmp_path)
    prs = Presentation(str(path))
    leaky = []
    for si, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if not shape.text_frame.text.strip():
                continue
            mode = shape.text_frame.auto_size
            if mode is not None and mode != MSO_AUTO_SIZE.NONE:
                leaky.append(f"slide {si+1}: shape auto_size={mode}")
    assert not leaky, (
        "Textboxes without auto_size=NONE will silently overflow:\n  "
        + "\n  ".join(leaky)
    )


# -------------------------------------------------------------- B2: anchor_word

def test_anchor_word_handles_cjk_title() -> None:
    """B2 regression: a CJK title like ``第三层:护栏`` used to produce empty
    output because the old logic stripped non-ASCII then took the longest
    word. Expect a short (≤ 6-char) CJK anchor instead."""
    from src.services.assistant.docgen.renderers.layouts.helpers import anchor_word

    # Prefix stripped, CJK tail returned
    assert anchor_word("第三层:护栏") == "护栏"
    assert anchor_word("第 3 层 — 护栏与防护") == "护栏与防护"
    # ASCII-preamble form
    assert anchor_word("Layer 3: 可观测性") == "可观测性"
    # Majority-CJK, no prefix
    assert anchor_word("架构设计总览") == "架构设计总览"
    # Long CJK gets capped at 6 chars
    out = anchor_word("非常非常非常长的中文标题文字")
    assert 1 <= len(out) <= 6

    # ASCII path still works for backward compat
    assert anchor_word("Layer 3: Guardrails") == "Guardrails"
    assert anchor_word("Observability") == "Observability"


# -------------------------------------------------------------- B3: split_head_tail

def test_split_head_tail_preserves_urls() -> None:
    """B3 regression: a label containing a URL used to mis-split on the
    ``:`` inside ``https://``. After the fix, URL-scheme colons are
    skipped."""
    from src.services.assistant.docgen.renderers.layouts.helpers import split_head_tail

    # URL inside tail — split on the LABEL colon, not the URL scheme
    assert split_head_tail("Docs: https://example.com") == (
        "Docs",
        "https://example.com",
    )

    # URL with port — the scheme colon AND the port colon must NOT split
    head, tail = split_head_tail("URL: https://foo.bar:8080/x")
    assert head == "URL"
    assert tail == "https://foo.bar:8080/x"

    # Bare URL only — nothing to split on, whole thing is head
    assert split_head_tail("https://example.com/path") == (
        "https://example.com/path",
        "",
    )

    # Dash separator still works
    assert split_head_tail("Head — tail description") == (
        "Head",
        "tail description",
    )

    # Plain label:value still splits
    assert split_head_tail("Key: value") == ("Key", "value")
